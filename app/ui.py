"""AutoApply – main window, debug console, settings dialog."""
import asyncio
import csv
import html
import os
import re
import shutil
import subprocess
import tempfile
import threading
import time
import webbrowser
from datetime import datetime

from PyQt6.QtCore import Qt, QTimer, QRect, pyqtSignal
from PyQt6.QtGui import (
    QBrush, QColor, QFont, QFontDatabase, QIcon, QLinearGradient,
    QPainter, QPen, QPixmap, QPalette,
)
from PyQt6.QtWidgets import (
    QApplication, QCheckBox, QComboBox, QDialog, QFileDialog,
    QFormLayout, QFrame, QHBoxLayout, QLabel, QLineEdit,
    QMainWindow, QMenu, QMessageBox, QProgressBar, QPushButton,
    QScrollArea, QSizePolicy, QSlider, QSplitter, QSystemTrayIcon,
    QTabWidget, QTextEdit, QVBoxLayout, QWidget,
)

from app import db, search_engine, ai_score_engine
from app.log import _ACTIVITY_LOG, _log_activity
from app.theme import (
    CAT, MONO_FONT, P, _btn, _card, _label, _mono, _pill,
    _pipe_qss, _section_header, _sep, _qss, score_col, strip_html,
)
from app.widgets import JobCard, MiniGraph, ViewToggle
from backend.database import init_db
from backend.services import llm_service


# ── search progress helper ───────────────────────────────────────────────────

def _search_progress_pct(phase: str) -> int:
    if not phase or phase in ("Ready", "IDLE", "Initialising…"):
        return 2
    if "BA API" in phase:
        m = re.search(r'(\d+)/(\d+)', phase)
        if m:
            return max(2, int(int(m.group(1)) / max(int(m.group(2)), 1) * 58))
        return 5
    if "Arbeitnow" in phase:
        return 62
    if "Filtering" in phase:
        return 76
    if "Saving" in phase:
        return 88
    if "Done" in phase:
        return 100
    if "Cancelled" in phase or "Error" in phase:
        return 0
    return 10


# ── description renderer ─────────────────────────────────────────────────────

_SECTION_ACCENTS = {
    "aufgaben":      P["indigo"],
    "tätigkeiten":   P["indigo"],
    "profil":        P["amber"],
    "anforderungen": P["amber"],
    "qualifikation": P["amber"],
    "bieten":        P["green"],
    "benefits":      P["green"],
    "kontakt":       P["text3"],
    "ansprechpartner": P["text3"],
}

def _format_desc_html(desc_plain: str, c: dict, title: str = "") -> str:
    """Render plain-text job description as structured HTML (sections, bullets, highlights)."""
    # Build keyword regex from job title words (> 3 chars)
    kw_pattern = None
    title_words = sorted(
        {w.lower() for w in re.split(r'\W+', title) if len(w) > 3},
        key=len, reverse=True,
    )
    if title_words:
        try:
            kw_pattern = re.compile(
                r'\b(' + '|'.join(re.escape(w) for w in title_words) + r')\b',
                re.IGNORECASE,
            )
        except re.error:
            pass

    def hl(escaped: str) -> str:
        if not kw_pattern:
            return escaped
        return kw_pattern.sub(
            f'<b style="color:{c["text"]};font-weight:600;">\\1</b>', escaped
        )

    def _section_accent(header_text: str) -> str:
        low = header_text.lower()
        for key, col in _SECTION_ACCENTS.items():
            if key in low:
                return col
        return c["indigo"]

    out: list[str] = []
    bullet_buf: list[str] = []
    first_paragraph = True

    def flush_bullets() -> None:
        if not bullet_buf:
            return
        items = "".join(
            f'<li style="margin:4px 0;line-height:1.65;">{hl(b)}</li>'
            for b in bullet_buf
        )
        out.append(
            f'<div style="margin:2px 0 14px 0;padding:8px 12px 8px 8px;'
            f'background:{c["card2"]};border-radius:8px;">'
            f'<ul style="margin:0;padding-left:22px;color:{c["text2"]};font-size:13px;">'
            f'{items}</ul></div>'
        )
        bullet_buf.clear()

    for line in desc_plain.split('\n'):
        s = line.strip()
        if not s:
            flush_bullets()
            continue

        if s.startswith('•'):
            bullet_buf.append(html.escape(s[1:].strip()))
            continue

        flush_bullets()

        # Section header: short line ending with ':', no mid-sentence punctuation
        is_header = (
            s.endswith(':')
            and len(s) < 65
            and s.count('.') == 0
            and s.count(',') == 0
        )
        if is_header:
            accent = _section_accent(s)
            out.append(
                f'<div style="margin:20px 0 7px 0;padding:4px 0 4px 12px;'
                f'border-left:3px solid {accent};'
                f'color:{c["text"]};font-size:13px;font-weight:700;'
                f'letter-spacing:0.2px;">{html.escape(s)}</div>'
            )
        else:
            size = "14px" if first_paragraph else "13px"
            out.append(
                f'<p style="margin:0 0 8px 0;color:{c["text2"]};'
                f'font-size:{size};line-height:1.75;">'
                f'{hl(html.escape(s))}</p>'
            )
            first_paragraph = False

    flush_bullets()
    return "".join(out)


# ── main window ───────────────────────────────────────────────────────────────

class MainWindow(QMainWindow):
    _ba_desc_ready  = pyqtSignal(int, str)       # job_id, description
    _ai_job_ready   = pyqtSignal(int, int, str)  # job_id, score, reason
    _bew_progress_sig = pyqtSignal(int, str, str)  # pct (-1=indeterminate), step, detail
    _bew_pdf_ready    = pyqtSignal(str)             # pdf path to render

    def __init__(self):
        super().__init__()
        self._ba_desc_ready.connect(self._ba_desc_done)
        self._ai_job_ready.connect(self._ai_job_done)
        self._bew_progress_sig.connect(self._on_bew_progress)
        self._bew_pdf_ready.connect(self._render_bew_pdf)
        self._ai_render_timer = QTimer(self); self._ai_render_timer.setSingleShot(True)
        self._ai_render_timer.timeout.connect(self._render_list)
        self.setWindowTitle("AutoApply")
        self.resize(1500, 900)
        self.setMinimumSize(1100, 680)

        self._jobs: list[dict] = []
        self._selected: dict | None = None
        self._cards: dict[int, JobCard] = {}
        self._filter_cat         = "All"
        self._view               = "all"
        self._show_dismissed     = False
        self._active_tab         = 0
        self._top_jobs_tick      = 0
        self._quick_new_only     = False
        self._quick_status       = ""
        self._last_dismissed_id: int | None = None
        self._debug_win          = None
        self._total_jobs_cache   = 0
        self._ai_idle_tick       = 0
        self._ai_unscored_cache  = 0

        # bewerbung state
        self._pptx_path             = ""
        self._bew_creating          = False
        self._bewerbung_paths: dict[int, str] = {}
        self._current_bew_job_id    = 0
        self._current_bew_out_path  = ""
        self._bew_out_dir           = ""

        self._build()
        self._run_auto_dismiss()
        self._load_jobs()

        self._poll_timer = QTimer(self)
        self._poll_timer.timeout.connect(self._poll_search)
        self._poll_timer.timeout.connect(self._poll_ai)
        self._poll_timer.start(400)

        self._auto_timer = QTimer(self)
        self._auto_timer.timeout.connect(self._auto_search_check)
        self._auto_timer.start(60_000)

        self._tray = QSystemTrayIcon(self)
        self._tray.setIcon(_make_app_icon())
        tray_menu = QMenu()
        tray_menu.addAction("Show", self.show)
        tray_menu.addAction("Quit", QApplication.quit)
        self._tray.setContextMenu(tray_menu)
        self._tray.activated.connect(lambda _: (self.show(), self.raise_()))
        self._tray.show()

        self._update_last_search_label()

    def closeEvent(self, ev):
        m = db.get_settings()["prefs"].get("ollama_model", "qwen2.5:14b")
        threading.Thread(target=llm_service.unload_model, args=(m,), daemon=True).start()
        ev.accept()

    # ── layout ────────────────────────────────────────────────────────────────

    def _build(self):
        central = QWidget()
        self.setCentralWidget(central)
        root = QVBoxLayout(central)
        root.setContentsMargins(0, 0, 0, 0)
        root.setSpacing(0)
        root.addWidget(self._build_topbar())
        self._splitter = QSplitter(Qt.Orientation.Horizontal)
        self._splitter.setHandleWidth(2)
        self._splitter.addWidget(self._build_sidebar())
        self._splitter.addWidget(self._build_detail())
        self._splitter.setSizes([540, 960])
        self._splitter.setCollapsible(0, False)
        self._splitter.setCollapsible(1, False)
        root.addWidget(self._splitter, 1)

    def _build_topbar(self) -> QWidget:
        bar = QWidget()
        bar.setFixedHeight(38)
        bar.setStyleSheet(f"QWidget {{ background: {P['card']}; border-bottom: 1px solid {P['border']}; }}")
        lay = QHBoxLayout(bar)
        lay.setContentsMargins(18, 0, 18, 0)
        lay.setSpacing(0)

        # Logo
        px = QPixmap(24, 24); px.fill(QColor(0, 0, 0, 0))
        lp = QPainter(px); lp.setRenderHint(QPainter.RenderHint.Antialiasing)
        g = QLinearGradient(0, 0, 0, 24)
        g.setColorAt(0, QColor("#1e1e50")); g.setColorAt(1, QColor("#08081a"))
        lp.setBrush(QBrush(g)); lp.setPen(QPen(QColor("#6d6df5"), 1.5))
        lp.drawRoundedRect(1, 1, 22, 22, 5, 5)
        lp.setPen(QColor("#eeeef8")); lp.setFont(QFont("Segoe UI Variable", 12, QFont.Weight.Bold))
        lp.drawText(px.rect(), Qt.AlignmentFlag.AlignCenter, "A"); lp.end()
        logo = QLabel(); logo.setFixedSize(24, 24); logo.setPixmap(px)
        logo.setStyleSheet("background: transparent;")
        lay.addWidget(logo)

        self._top_dot = QLabel(" ●")
        self._top_dot.setStyleSheet(f"color: {P['green']}; font-size: 10px; background: transparent;")
        brand = QLabel("AUTOAPPLY")
        brand.setStyleSheet(
            f"color: {P['text']}; font-size: 11px; font-weight: 700; "
            f"letter-spacing: 3px; font-family: '{MONO_FONT}'; background: transparent;"
        )
        lay.addWidget(self._top_dot)
        lay.addWidget(brand)
        lay.addStretch()

        sep = QLabel("  |  ")
        sep.setStyleSheet(f"color: {P['border2']}; background: transparent;")
        self._top_status = _mono("IDLE",       color=P['text3'])
        self._top_jobs   = _mono("DB  — jobs", color=P['text3'])
        self._top_time   = _mono("",           color=P['text3'])
        for w in (self._top_status, sep, self._top_jobs, self._top_time):
            lay.addWidget(w); lay.addSpacing(16)
        return bar

    def _build_sidebar(self) -> QWidget:
        sb = QWidget()
        sb.setMinimumWidth(380); sb.setMaximumWidth(520)
        sb.setStyleSheet(f"background: {P['sidebar']};")
        lay = QVBoxLayout(sb)
        lay.setContentsMargins(14, 16, 14, 12)
        lay.setSpacing(0)

        # ── Search ────────────────────────────────────────────────────────────
        lay.addWidget(_section_header("SEARCH ENGINE", P['indigo']))
        lay.addSpacing(5)
        self._btn_search = _btn("Start Search", P['indigo'], P['indigo_d'], height=34, font_size=13)
        self._btn_search.clicked.connect(self._start_search)
        lay.addWidget(self._btn_search)
        lay.addSpacing(4)

        self._progress = QProgressBar()
        self._progress.setFixedHeight(26)
        self._progress.setRange(0, 100); self._progress.setValue(0)
        self._progress.setTextVisible(True); self._progress.setFormat("  READY")
        self._btn_cancel = QPushButton("✕")
        self._btn_cancel.setFixedSize(26, 26)
        self._btn_cancel.setStyleSheet(
            f"QPushButton {{ background: transparent; color: {P['red']}; "
            f"border: 1px solid {P['red_bg']}; border-radius: 5px; "
            f"font-size: 11px; font-weight: 700; padding: 0; }}"
            f"QPushButton:hover {{ background: {P['red_bg']}; }}"
        )
        self._btn_cancel.setEnabled(False)
        self._btn_cancel.clicked.connect(self._cancel_search)
        prog_row = QHBoxLayout(); prog_row.setSpacing(4)
        prog_row.addWidget(self._progress, 1); prog_row.addWidget(self._btn_cancel)
        lay.addLayout(prog_row); lay.addSpacing(3)

        phase_row = QHBoxLayout(); phase_row.setSpacing(0)
        self._lbl_phase = _mono("IDLE", size=10, color=P['text3'])
        self._lbl_search_meta = _mono("", size=10, color=P['text3'])
        self._lbl_search_meta.setAlignment(Qt.AlignmentFlag.AlignRight)
        phase_row.addWidget(self._lbl_phase, 1); phase_row.addWidget(self._lbl_search_meta)
        lay.addLayout(phase_row); lay.addSpacing(2)
        self._lbl_last_search = _mono("Last search: never", size=10, color=P['text3'])
        lay.addWidget(self._lbl_last_search); lay.addSpacing(5)

        # stats strip (visible during search, hidden when idle)
        self._stats_frame = QFrame(); self._stats_frame.setObjectName("sf")
        self._stats_frame.setStyleSheet(
            f"QFrame#sf {{ background: {P['card2']}; border-radius: 8px; border: 1px solid {P['border']}; }}"
        )
        sl = QHBoxLayout(self._stats_frame); sl.setContentsMargins(0, 4, 0, 4)
        self._st_fetched  = self._stat_cell(sl, "FETCHED", "—")
        self._stat_vsep(sl)
        self._st_filtered = self._stat_cell(sl, "FILTERED", "—")
        self._stat_vsep(sl)
        self._st_saved    = self._stat_cell(sl, "NEW", "—")
        self._stat_vsep(sl)
        self._st_total    = self._stat_cell(sl, "TOTAL DB", "—")
        lay.addWidget(self._stats_frame)
        self._lbl_idle_db = _mono("", size=11, color=P['text3'])
        lay.addWidget(self._lbl_idle_db)
        self._lbl_idle_db.hide()
        lay.addSpacing(4)

        self._lbl_funnel = QLabel()
        self._lbl_funnel.setStyleSheet("background: transparent;")
        self._lbl_funnel.setTextFormat(Qt.TextFormat.RichText)
        self._lbl_funnel.setText(self._funnel_html({"new": 0, "applied": 0, "interview": 0, "offer": 0}))
        lay.addWidget(self._lbl_funnel); lay.addSpacing(4)

        lay.addSpacing(6)
        lay.addWidget(_sep()); lay.addSpacing(6)

        # ── AI Scorer ─────────────────────────────────────────────────────────
        lay.addWidget(_section_header("AI SCORER", P['purple']))
        lay.addSpacing(4)

        self._ai_progress = QProgressBar()
        self._ai_progress.setFixedHeight(22); self._ai_progress.setRange(0, 100)
        self._ai_progress.setValue(0); self._ai_progress.setTextVisible(True)
        self._ai_progress.setFormat("  IDLE")
        self._ai_progress.setStyleSheet(
            f"QProgressBar {{ background: {P['card2']}; border: 1px solid {P['border2']}; "
            f"border-radius: 5px; height: 22px; text-align: center; color: {P['text']}; "
            f"font-size: 11px; font-weight: 700; font-family: '{MONO_FONT}'; }}"
            f"QProgressBar::chunk {{ border-radius: 4px; background: qlineargradient("
            f"x1:0,y1:0,x2:1,y2:0, stop:0 #4b0082, stop:0.5 #7b2fbe, stop:1 {P['purple']}); }}"
        )
        lay.addWidget(self._ai_progress); lay.addSpacing(3)

        ai_info_row = QHBoxLayout(); ai_info_row.setSpacing(0)
        self._lbl_ai_phase = _mono("idle", size=10, color=P['text3'])
        self._lbl_ai_eta   = _mono("", size=10, color=P['text3'])
        self._lbl_ai_eta.setAlignment(Qt.AlignmentFlag.AlignRight)
        ai_info_row.addWidget(self._lbl_ai_phase, 1); ai_info_row.addWidget(self._lbl_ai_eta)
        lay.addLayout(ai_info_row); lay.addSpacing(4)

        self._ai_auto_enabled = False
        self._btn_ai_toggle = QPushButton("AUTO  OFF")
        self._btn_ai_toggle.setFixedHeight(26)
        self._btn_ai_toggle.setCheckable(True)
        self._btn_ai_toggle.setChecked(False)
        self._btn_ai_toggle.clicked.connect(self._on_ai_toggle)
        self._apply_ai_toggle_style()

        ai_btn_row = QHBoxLayout(); ai_btn_row.setSpacing(4)
        self._btn_ai_start  = _btn("✦  Score now", P['purple_bg'], "#2a0e50",
                                   height=26, font_size=11, color=P['purple'])
        self._btn_ai_cancel = QPushButton("✕ Cancel")
        self._btn_ai_cancel.setFixedHeight(26)
        self._btn_ai_cancel.setStyleSheet(
            f"QPushButton {{ background: {P['red_bg']}; color: {P['red']}; "
            f"border: 1px solid {P['red']}44; border-radius: 5px; font-size: 11px; font-weight: 700; padding: 0 10px; }}"
            f"QPushButton:hover {{ background: {P['red']}33; border-color: {P['red']}; }}"
            f"QPushButton:disabled {{ background: transparent; color: {P['text3']}; border-color: {P['border']}; }}"
        )
        self._btn_ai_cancel.setEnabled(False)
        self._btn_ai_start.clicked.connect(self._start_ai_scoring)
        def _on_cancel():
            ai_score_engine.cancel_scoring()
            self._slbl(self._lbl_ai_phase, "CANCELLING…", P['red'])
            self._btn_ai_cancel.setEnabled(False)
        self._btn_ai_cancel.clicked.connect(_on_cancel)
        ai_btn_row.addWidget(self._btn_ai_toggle)
        ai_btn_row.addWidget(self._btn_ai_start, 1)
        ai_btn_row.addWidget(self._btn_ai_cancel)
        lay.addLayout(ai_btn_row); lay.addSpacing(6)
        lay.addWidget(_sep()); lay.addSpacing(6)

        # ── Filter ────────────────────────────────────────────────────────────
        lay.addWidget(_section_header("FILTER & VIEW", P['amber']))
        lay.addSpacing(4)
        self._view_toggle = ViewToggle(["All", "Saved", "Applied"])
        self._view_toggle.changed.connect(self._on_view_change)
        lay.addWidget(self._view_toggle); lay.addSpacing(4)

        # quick-filter chips
        qf_row = QHBoxLayout(); qf_row.setSpacing(4)
        def _qchip(label, tip=""):
            b = QPushButton(label); b.setCheckable(True); b.setFixedHeight(24)
            b.setToolTip(tip)
            b.setStyleSheet(
                f"QPushButton{{background:{P['card2']};color:{P['text3']};border-radius:5px;"
                f"font-size:10px;font-weight:600;padding:0 8px;}}"
                f"QPushButton:checked{{background:{P['indigo_bg']};color:{P['indigo']};}}"
                f"QPushButton:hover{{background:{P['card3']};}}"
            )
            return b
        self._chip_new       = _qchip("● Neu",      "Nur heute gefundene Jobs")
        self._chip_interview = _qchip("Interview",  "Nur Jobs im Interview-Status")
        self._chip_new.clicked.connect(self._on_chip_new)
        self._chip_interview.clicked.connect(self._on_chip_interview)
        qf_row.addWidget(self._chip_new)
        qf_row.addWidget(self._chip_interview)
        qf_row.addStretch()
        lay.addLayout(qf_row); lay.addSpacing(4)

        self._search_input = QLineEdit()
        self._search_input.setPlaceholderText("Search title or company…")
        self._search_input.setFixedHeight(30)
        self._search_input.textChanged.connect(self._on_filter)
        lay.addWidget(self._search_input); lay.addSpacing(4)

        filter_row = QHBoxLayout(); filter_row.setSpacing(4)
        self._cat_combo = QComboBox()
        self._cat_combo.addItems(["All Categories", "IT", "Wirtschaft", "Unknown"])
        self._cat_combo.currentTextChanged.connect(self._on_filter)
        self._sort_combo = QComboBox()
        self._sort_combo.addItem("Score ↓",   "score")
        self._sort_combo.addItem("Date ↓",    "date")
        self._sort_combo.addItem("Company ↑", "company")
        self._sort_combo.currentIndexChanged.connect(self._on_filter)
        filter_row.addWidget(self._cat_combo, 3); filter_row.addWidget(self._sort_combo, 2)
        lay.addLayout(filter_row); lay.addSpacing(4)

        slider_row = QHBoxLayout(); slider_row.setSpacing(6)
        self._lbl_min_score = _mono("MIN SCORE  0", size=10, color=P['text3'])
        self._score_slider  = QSlider(Qt.Orientation.Horizontal)
        self._score_slider.setRange(0, 100); self._score_slider.setValue(0)
        self._score_slider.setFixedHeight(16)
        self._score_slider.valueChanged.connect(self._on_score_slider)
        slider_row.addWidget(self._lbl_min_score)
        slider_row.addWidget(self._score_slider, 1)
        lay.addLayout(slider_row); lay.addSpacing(3)

        _chk_style = (
            f"QCheckBox {{ color: {P['text3']}; font-size: 11px; spacing: 5px; }}"
            f"QCheckBox::indicator {{ width: 13px; height: 13px; border-radius: 3px; "
            f"border: 1px solid {P['border2']}; background: {P['card2']}; }}"
            f"QCheckBox::indicator:checked {{ background: {P['indigo']}; border-color: {P['indigo']}; }}"
        )
        chk_row = QHBoxLayout()
        self._chk_dismissed = QCheckBox("Dismissed")
        self._chk_dismissed.setStyleSheet(_chk_style)
        self._chk_dismissed.stateChanged.connect(self._on_filter)
        self._chk_ai_only = QCheckBox("✦ AI only")
        self._chk_ai_only.setStyleSheet(_chk_style)
        self._chk_ai_only.stateChanged.connect(self._on_filter)
        self._lbl_count = _mono("0 jobs", color=P['text3'])
        self._lbl_count.setAlignment(Qt.AlignmentFlag.AlignRight)
        chk_row.addWidget(self._chk_dismissed)
        chk_row.addSpacing(10)
        chk_row.addWidget(self._chk_ai_only)
        chk_row.addStretch(); chk_row.addWidget(self._lbl_count)
        lay.addLayout(chk_row); lay.addSpacing(6)
        lay.addWidget(_sep()); lay.addSpacing(4)

        # ── Job list ──────────────────────────────────────────────────────────
        lay.addSpacing(2)
        self._scroll = QScrollArea()
        self._scroll.setWidgetResizable(True)
        self._scroll.setHorizontalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAlwaysOff)
        self._scroll.setStyleSheet("background: transparent;")
        self._list_widget  = QWidget(); self._list_widget.setStyleSheet("background: transparent;")
        self._list_layout  = QVBoxLayout(self._list_widget)
        self._list_layout.setContentsMargins(0, 0, 6, 0); self._list_layout.setSpacing(5)
        self._list_layout.addStretch()
        self._scroll.setWidget(self._list_widget)
        lay.addWidget(self._scroll, 1); lay.addSpacing(10)
        lay.addWidget(_sep()); lay.addSpacing(8)

        bot = QHBoxLayout(); bot.setSpacing(4)
        self._btn_settings    = _btn("⚙",          P['card2'],   P['card3'],    height=28, font_size=12, fixed_width=32)
        self._btn_settings.setToolTip("Settings")
        self._btn_stats       = _btn("📊",          P['card2'],   P['card3'],    height=28, font_size=12, fixed_width=32)
        self._btn_stats.setToolTip("Statistics")
        self._btn_dismiss_all = _btn("✕ Dismiss",   P['card2'],   P['card3'],    height=28, font_size=11)
        self._btn_dismiss_all.setToolTip("Dismiss all currently visible jobs")
        self._btn_more        = _btn("⋮",           P['card2'],   P['card3'],    height=28, font_size=14, fixed_width=32)
        self._btn_more.setToolTip("More options")
        def _show_more_menu():
            m = QMenu(self)
            m.setStyleSheet(
                f"QMenu {{ background: {P['card2']}; color: {P['text']}; "
                f"border: 1px solid {P['border2']}; border-radius: 6px; padding: 4px; }}"
                f"QMenu::item {{ padding: 6px 18px; border-radius: 4px; }}"
                f"QMenu::item:selected {{ background: {P['card3']}; }}"
                f"QMenu::separator {{ background: {P['border']}; height: 1px; margin: 4px 0; }}"
            )
            m.addAction("⬡  Debug Console", self._open_debug)
            m.addSeparator()
            m.addAction("↓  Export CSV", self._export_csv)
            m.addSeparator()
            m.addAction("🗑  Clear all jobs", self._clear_jobs)
            m.exec(self._btn_more.mapToGlobal(self._btn_more.rect().bottomLeft()))
        self._btn_settings.clicked.connect(self._open_settings)
        self._btn_stats.clicked.connect(self._open_stats)
        self._btn_dismiss_all.clicked.connect(self._dismiss_visible)
        self._btn_more.clicked.connect(_show_more_menu)
        bot.addWidget(self._btn_settings)
        bot.addWidget(self._btn_stats)
        bot.addWidget(self._btn_dismiss_all, 1)
        bot.addWidget(self._btn_more)
        lay.addLayout(bot)
        return sb

    def _stat_cell(self, lay: QHBoxLayout, label: str, value: str) -> QLabel:
        cell = QWidget(); cell.setStyleSheet("background: transparent;")
        vl = QVBoxLayout(cell); vl.setContentsMargins(0, 0, 0, 0); vl.setSpacing(2)
        val = QLabel(value)
        val.setStyleSheet(
            f"color: {P['text']}; font-size: 20px; font-weight: 700; "
            f"font-family: '{MONO_FONT}'; background: transparent;"
        )
        val.setAlignment(Qt.AlignmentFlag.AlignCenter)
        lbl = _mono(label, size=11, color=P['text3'])
        lbl.setAlignment(Qt.AlignmentFlag.AlignCenter)
        vl.addWidget(val); vl.addWidget(lbl)
        lay.addWidget(cell, 1)
        return val

    def _stat_vsep(self, lay: QHBoxLayout):
        line = QFrame(); line.setFrameShape(QFrame.Shape.VLine)
        line.setFixedWidth(1); line.setStyleSheet(f"background: {P['border']}; border: none;")
        lay.addWidget(line)

    def _build_detail(self) -> QWidget:
        panel = QWidget()
        lay = QVBoxLayout(panel)
        lay.setContentsMargins(16, 14, 16, 14); lay.setSpacing(8)

        # ── header card ───────────────────────────────────────────────────────
        hdr = _card(radius=14)
        hl = QVBoxLayout(hdr); hl.setContentsMargins(20, 16, 20, 14); hl.setSpacing(4)

        tr = QHBoxLayout(); tr.setSpacing(12)
        self._lbl_title = _label("Select a job from the list", size=21, color=P['text'], bold=True)
        self._lbl_score = QLabel("—")
        self._lbl_score.setStyleSheet(
            f"background: {P['border']}; color: {P['text3']}; border-radius: 8px; "
            f"padding: 5px 14px; font-size: 16px; font-weight: 700; font-family: '{MONO_FONT}';"
        )
        self._lbl_score.setSizePolicy(QSizePolicy.Policy.Fixed, QSizePolicy.Policy.Fixed)
        self._lbl_score.setAlignment(Qt.AlignmentFlag.AlignCenter)
        tr.addWidget(self._lbl_title, 1); tr.addWidget(self._lbl_score, 0, Qt.AlignmentFlag.AlignTop)
        hl.addLayout(tr)

        self._lbl_meta   = _label("", size=11, color=P['text2'])
        self._lbl_reason = _mono("", size=10, color=P['text3'])
        hl.addWidget(self._lbl_meta); hl.addWidget(self._lbl_reason)

        self._tag_row = QHBoxLayout(); self._tag_row.setSpacing(6)
        self._tag_row.setAlignment(Qt.AlignmentFlag.AlignLeft)
        hl.addLayout(self._tag_row)
        hl.addSpacing(6); hl.addWidget(_sep()); hl.addSpacing(6)

        # action buttons
        ab = QHBoxLayout(); ab.setSpacing(8)
        self._btn_open     = _btn("Open Job",    P['indigo'],    P['indigo_d'], height=34)
        self._btn_save     = _btn("Save",        P['card2'],     P['card3'],    height=34)
        self._btn_dismiss  = _btn("Dismiss",     P['card2'],     P['card3'],    height=34)
        self._btn_undo_dis = _btn("↩",           P['amber_bg'],  P['amber_bg'], height=34, fixed_width=40)
        self._btn_undo_dis.setToolTip("Undo last dismiss")
        self._btn_undo_dis.setEnabled(False)
        self._btn_rescore  = _btn("✦ Re-score",  P['purple_bg'], "#2a0e50",    height=34, color=P['purple'])
        self._btn_bew_go   = _btn("📄 Bewerbung", P['card2'],    P['card3'],    height=34)
        self._btn_bew_go.setToolTip("Switch to Bewerbung tab and create PDF  [B]")
        for b in (self._btn_open, self._btn_save, self._btn_dismiss,
                  self._btn_undo_dis, self._btn_rescore, self._btn_bew_go):
            ab.addWidget(b)
        ab.addStretch()
        _btn_keys = QPushButton("⌨")
        _btn_keys.setFixedSize(34, 34)
        _btn_keys.setToolTip("Keyboard shortcuts")
        _btn_keys.setStyleSheet(
            f"QPushButton {{ background: transparent; color: {P['text3']}; border: none; "
            f"font-size: 16px; }}"
            f"QPushButton:hover {{ color: {P['text2']}; }}"
        )
        def _show_shortcuts():
            msg = QMessageBox(self)
            msg.setWindowTitle("Keyboard Shortcuts")
            msg.setText(
                "<table cellspacing='8'>"
                "<tr><td><b>S</b></td><td>Save / Unsave</td></tr>"
                "<tr><td><b>D</b></td><td>Dismiss</td></tr>"
                "<tr><td><b>A</b></td><td>Mark as Applied</td></tr>"
                "<tr><td><b>B</b></td><td>Bewerbung erstellen (PDF)</td></tr>"
                "<tr><td><b>N</b></td><td>Notes tab</td></tr>"
                "<tr><td><b>O</b></td><td>Open job in browser</td></tr>"
                "<tr><td><b>↑ ↓</b></td><td>Navigate job list</td></tr>"
                "<tr><td><b>Tab</b></td><td>Nächster ungelesener Job</td></tr>"
                "</table>"
            )
            msg.setStyleSheet(f"QMessageBox {{ background: {P['bg']}; }} QLabel {{ color: {P['text']}; }}")
            msg.exec()
        _btn_keys.clicked.connect(_show_shortcuts)
        ab.addWidget(_btn_keys)
        hl.addLayout(ab)

        # status pipeline (Applied / Interview / Offer / Rejected)
        pipe = QHBoxLayout(); pipe.setSpacing(6)
        pipe.addWidget(_mono("STATUS", size=10, color=P['text3'])); pipe.addSpacing(4)
        self._btn_applied_pipe = QPushButton("Applied")
        self._btn_interview    = QPushButton("Interview")
        self._btn_offer        = QPushButton("Offer")
        self._btn_rejected     = QPushButton("Rejected")
        for btn, col in [(self._btn_applied_pipe, P['indigo']),
                         (self._btn_interview,    P['amber']),
                         (self._btn_offer,        P['green']),
                         (self._btn_rejected,     P['red'])]:
            btn.setFixedHeight(32)
            btn.setStyleSheet(_pipe_qss(col, False))
        self._btn_applied_pipe.clicked.connect(lambda: self._set_status_pipeline("applied"))
        self._btn_interview.clicked.connect(lambda: self._set_status_pipeline("interview"))
        self._btn_offer.clicked.connect(lambda: self._set_status_pipeline("offer"))
        self._btn_rejected.clicked.connect(lambda: self._set_status_pipeline("rejected"))
        for b in (self._btn_applied_pipe, self._btn_interview, self._btn_offer, self._btn_rejected):
            pipe.addWidget(b)
        pipe.addStretch()
        hl.addLayout(pipe)
        lay.addWidget(hdr)

        self._btn_open.clicked.connect(self._open_job)
        self._btn_save.clicked.connect(self._toggle_save)
        self._btn_dismiss.clicked.connect(self._toggle_dismiss)
        self._btn_undo_dis.clicked.connect(self._undo_dismiss)
        self._btn_rescore.clicked.connect(self._rescore_current)
        self._btn_bew_go.clicked.connect(self._go_to_bewerbung)

        # ── tabs ──────────────────────────────────────────────────────────────
        self._tabs = QTabWidget()
        self._tabs.setDocumentMode(True)
        self._tabs.currentChanged.connect(lambda i: setattr(self, "_active_tab", i))

        # Job Profile tab
        desc_w = QWidget(); desc_w.setStyleSheet(f"background: {P['card']};")
        desc_l = QVBoxLayout(desc_w); desc_l.setContentsMargins(0, 0, 0, 0)
        self._txt_desc = QTextEdit()
        self._txt_desc.setReadOnly(True)
        self._txt_desc.setStyleSheet("border: none; padding: 14px;")
        _set_palette(self._txt_desc, P['card'], P['text2'])
        desc_l.addWidget(self._txt_desc, 1)
        self._tabs.addTab(desc_w, "Job Profile")

        # Bewerbung tab (PDF viewer)
        self._tabs.addTab(self._build_bewerbung_tab(), "Bewerbung")

        # Notes tab
        notes_w = QWidget(); notes_w.setStyleSheet(f"background: {P['card']};")
        notes_l = QVBoxLayout(notes_w); notes_l.setContentsMargins(16, 14, 16, 14); notes_l.setSpacing(8)
        notes_hdr = QHBoxLayout()
        notes_hdr.addWidget(_label("Notes", size=11, color=P['text2'], bold=True))
        notes_hdr.addStretch()
        self._lbl_notes_saved = _label("", size=9, color=P['text3'])
        notes_hdr.addWidget(self._lbl_notes_saved)
        notes_l.addLayout(notes_hdr)
        self._txt_notes = QTextEdit()
        self._txt_notes.setStyleSheet("border: none;")
        _set_palette(self._txt_notes, P['card'], P['text'])
        self._txt_notes.setPlaceholderText("Your notes for this job… (auto-saved)")
        notes_l.addWidget(self._txt_notes, 1)
        self._tabs.addTab(notes_w, "Notes")
        self._notes_timer = QTimer(self); self._notes_timer.setSingleShot(True)
        self._notes_timer.timeout.connect(self._save_notes)
        self._txt_notes.textChanged.connect(self._on_notes_changed)

        lay.addWidget(self._tabs, 1)
        return panel

    # ── Bewerbung PPTX tab ────────────────────────────────────────────────────

    def _build_bewerbung_tab(self) -> QWidget:
        w = QWidget(); w.setStyleSheet(f"background: {P['card']};")
        lay = QVBoxLayout(w); lay.setContentsMargins(16, 14, 16, 14); lay.setSpacing(8)

        # ── top bar ──────────────────────────────────────────────────────────
        bar = QHBoxLayout(); bar.setSpacing(6)
        self._lbl_pptx_path = _label("No template loaded", size=11, color=P['text3'])
        self._lbl_pptx_path.setWordWrap(False)
        btn_browse      = _btn("📂 Load Template", P['card2'],   P['card3'],    height=32, font_size=11)
        self._btn_create = _btn("📄 Create PDF",   P['indigo'],  P['indigo_d'], height=32, font_size=11)
        bar.addWidget(self._lbl_pptx_path, 1)
        bar.addWidget(btn_browse)
        bar.addWidget(self._btn_create)
        lay.addLayout(bar); lay.addWidget(_sep())

        # ── progress panel ────────────────────────────────────────────────────
        prog_frame = QFrame(); prog_frame.setObjectName("bpf")
        prog_frame.setStyleSheet(
            f"QFrame#bpf {{ background: {P['card2']}; border-radius: 10px; "
            f"border: 1px solid {P['border']}; }}"
        )
        pfl = QVBoxLayout(prog_frame); pfl.setContentsMargins(16, 12, 16, 12); pfl.setSpacing(5)

        step_row = QHBoxLayout(); step_row.setSpacing(0)
        self._lbl_bew_step = _label("Ready", size=12, color=P['text'], bold=True)
        self._lbl_bew_pct  = _mono("", size=10, color=P['text3'])
        self._lbl_bew_pct.setAlignment(Qt.AlignmentFlag.AlignRight)
        step_row.addWidget(self._lbl_bew_step, 1); step_row.addWidget(self._lbl_bew_pct)
        pfl.addLayout(step_row)

        self._bew_progress = QProgressBar()
        self._bew_progress.setFixedHeight(6)
        self._bew_progress.setRange(0, 100); self._bew_progress.setValue(0)
        self._bew_progress.setTextVisible(False)
        self._bew_progress.setStyleSheet(
            f"QProgressBar {{ background: {P['border']}; border-radius: 3px; border: none; }}"
            f"QProgressBar::chunk {{ border-radius: 3px; background: qlineargradient("
            f"x1:0,y1:0,x2:1,y2:0, stop:0 {P['indigo']}, stop:1 {P['purple']}); }}"
        )
        pfl.addWidget(self._bew_progress)

        self._lbl_bew_detail = _label("Select a job and press Create PDF", size=10, color=P['text3'])
        pfl.addWidget(self._lbl_bew_detail)

        self._lbl_bew_target = _label("", size=11, color=P['text2'], bold=True)
        self._lbl_bew_target.setWordWrap(True)
        pfl.addWidget(self._lbl_bew_target)
        lay.addWidget(prog_frame)

        # ── post-creation actions ─────────────────────────────────────────────
        self._bew_action_row = QWidget()
        act_lay = QHBoxLayout(self._bew_action_row)
        act_lay.setContentsMargins(0, 2, 0, 2); act_lay.setSpacing(6)
        self._lbl_bew_job_status = _label("", size=10, color=P['text3'])
        self._btn_open_bew_folder  = _btn("📂 Ordner öffnen",       P['card2'], P['card3'], height=28, font_size=10)
        self._btn_mark_applied_bew = _btn("✓ Als beworben markieren", P['green'], P['green'], height=28, font_size=10)
        act_lay.addWidget(self._lbl_bew_job_status, 1)
        act_lay.addWidget(self._btn_open_bew_folder)
        act_lay.addWidget(self._btn_mark_applied_bew)
        self._btn_open_bew_folder.hide()
        self._btn_mark_applied_bew.hide()
        lay.addWidget(self._bew_action_row)
        lay.addWidget(_sep())

        def _open_bew_folder():
            folder = self._bew_out_dir
            if not folder and self._pptx_path:
                folder = os.path.join(os.path.dirname(self._pptx_path), "Bewerbungen")
            if folder:
                os.makedirs(folder, exist_ok=True)
                os.startfile(folder)

        def _mark_applied_from_bew():
            jid = self._current_bew_job_id
            if jid:
                db.set_status(jid, "applied")
                if jid in self._cards:
                    self._cards[jid].job["status"] = "applied"
                    self._cards[jid]._rebuild_tags()
                    self._cards[jid]._apply_style()
                self._btn_mark_applied_bew.setEnabled(False)
                self._btn_mark_applied_bew.setText("✓ Beworben")

        self._btn_open_bew_folder.clicked.connect(_open_bew_folder)
        self._btn_mark_applied_bew.clicked.connect(_mark_applied_from_bew)

        # ── PDF preview scroll area ───────────────────────────────────────────
        self._bew_scroll = QScrollArea()
        self._bew_scroll.setWidgetResizable(True)
        self._bew_scroll.setHorizontalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAsNeeded)
        self._bew_scroll.setStyleSheet("background: transparent; border: none;")
        self._bew_page_container = QWidget()
        self._bew_page_container.setStyleSheet("background: transparent;")
        self._bew_page_layout = QVBoxLayout(self._bew_page_container)
        self._bew_page_layout.setContentsMargins(0, 0, 0, 0); self._bew_page_layout.setSpacing(8)
        self._bew_page_layout.addStretch()
        self._bew_scroll.setWidget(self._bew_page_container)
        lay.addWidget(self._bew_scroll, 1)

        # ── load saved PPTX path ─────────────────────────────────────────────
        saved = db.get_settings()["prefs"].get("cv_pptx_path", "")
        if saved and os.path.isfile(saved):
            self._pptx_path = saved
            self._lbl_pptx_path.setText(os.path.basename(saved))
            QTimer.singleShot(300, lambda: self._preview_template(saved))

        def _browse():
            path, _ = QFileDialog.getOpenFileName(self, "Select Bewerbung PPTX", "", "PowerPoint (*.pptx)")
            if path:
                self._pptx_path = path
                self._lbl_pptx_path.setText(os.path.basename(path))
                s = db.get_settings(); s["prefs"]["cv_pptx_path"] = path
                db.save_settings(s["profile"], s["prefs"])
                self._preview_template(path)

        btn_browse.clicked.connect(_browse)
        self._btn_create.clicked.connect(self._create_application_pdf)
        return w

    def _update_bew_status(self):
        """Update status label in Bewerbung tab for currently selected job."""
        job = getattr(self, "_selected", None)
        if not job:
            self._lbl_bew_job_status.setText("")
            self._lbl_bew_target.setText("")
            return
        company = job.get("company") or ""
        target_text = job["title"] + (f"  ·  {company}" if company else "")
        self._lbl_bew_target.setText(target_text)
        jid = job["id"]
        path = self._bewerbung_paths.get(jid, "")
        if path:
            self._lbl_bew_job_status.setText(f"✔ {os.path.basename(path)}")
            self._lbl_bew_job_status.setStyleSheet(f"color: {P['green']}; font-size: 10px;")
            self._btn_open_bew_folder.show()
            self._btn_mark_applied_bew.show()
            already = job.get("status") == "applied"
            self._btn_mark_applied_bew.setEnabled(not already)
            self._btn_mark_applied_bew.setText("✓ Beworben" if already else "✓ Als beworben markieren")
        else:
            self._lbl_bew_job_status.setText("Noch keine Bewerbung erstellt")
            self._lbl_bew_job_status.setStyleSheet(f"color: {P['text3']}; font-size: 10px;")
            self._btn_open_bew_folder.hide()
            self._btn_mark_applied_bew.hide()

    def _on_bew_progress(self, pct: int, step: str, detail: str):
        if pct < 0:
            self._bew_progress.setRange(0, 0)  # pulsing indeterminate
        else:
            self._bew_progress.setRange(0, 100)
            self._bew_progress.setValue(pct)
        self._lbl_bew_step.setText(step)
        self._lbl_bew_detail.setText(detail)
        if pct == 100:
            self._lbl_bew_pct.setText("100 %")
            self._btn_create.setEnabled(True)
            if self._bew_creating and self._current_bew_out_path:
                jid = self._current_bew_job_id
                self._bewerbung_paths[jid] = self._current_bew_out_path
                self._update_bew_status()
            self._bew_creating = False
        elif pct > 0:
            self._lbl_bew_pct.setText(f"{pct} %")
        elif pct == 0:
            self._lbl_bew_pct.setText("")
            self._btn_create.setEnabled(True)
            self._bew_creating = False
        else:
            self._lbl_bew_pct.setText("…")

    def _render_bew_pdf(self, pdf_path: str):
        import fitz
        while self._bew_page_layout.count() > 1:
            item = self._bew_page_layout.takeAt(0)
            if item.widget(): item.widget().deleteLater()
        try:
            doc = fitz.open(pdf_path)
            for page in doc:
                pix  = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
                qpix = QPixmap(); qpix.loadFromData(pix.tobytes("png"))
                lbl  = QLabel(); lbl.setPixmap(qpix); lbl.setFixedSize(qpix.size())
                lbl.setStyleSheet("background: white; border-radius: 4px;")
                wrapper = QWidget(); wrapper.setStyleSheet("background: transparent;")
                wl = QHBoxLayout(wrapper); wl.setContentsMargins(0, 0, 0, 0)
                wl.addStretch(); wl.addWidget(lbl); wl.addStretch()
                self._bew_page_layout.insertWidget(self._bew_page_layout.count() - 1, wrapper)
            doc.close()
        except Exception as exc:
            self._lbl_bew_detail.setText(f"Preview error: {exc}")

    def _find_soffice(self) -> str | None:
        for p in [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            "soffice",
        ]:
            if os.path.isfile(p) or shutil.which(p):
                return p
        return None

    def _preview_template(self, pptx_path: str):
        self._bew_progress_sig.emit(-1, "Loading template preview…", os.path.basename(pptx_path))
        def _work():
            soffice = self._find_soffice()
            if not soffice:
                self._bew_progress_sig.emit(0, "Preview unavailable", "LibreOffice not found"); return
            with tempfile.TemporaryDirectory() as tmp:
                result = subprocess.run(
                    [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp, pptx_path],
                    capture_output=True, timeout=60,
                )
                base    = os.path.splitext(os.path.basename(pptx_path))[0]
                tmp_pdf = os.path.join(tmp, base + ".pdf")
                if result.returncode == 0 and os.path.isfile(tmp_pdf):
                    fd, dest = tempfile.mkstemp(suffix=".pdf", prefix="bew_preview_")
                    os.close(fd)
                    shutil.copy2(tmp_pdf, dest)
                    self._bew_progress_sig.emit(100, "Template loaded", os.path.basename(pptx_path))
                    self._bew_pdf_ready.emit(dest)
                else:
                    self._bew_progress_sig.emit(0, "Preview error", "Could not convert template to PDF")
        threading.Thread(target=_work, daemon=True).start()

    def _create_application_pdf(self):
        if self._bew_creating:
            return
        job = self._selected
        if not job:
            QMessageBox.warning(self, "No Job Selected", "Please select a job first."); return
        if not self._pptx_path or not os.path.isfile(self._pptx_path):
            QMessageBox.warning(self, "No Template", "Please load a PPTX template first."); return

        title   = job.get("title", "").upper()
        today   = datetime.now().strftime("%d.%m.%Y")
        company = (job.get("company") or "").replace("/", "-").replace("\\", "-") or "Bewerbung"

        out_dir = os.path.join(os.path.dirname(self._pptx_path), "Bewerbungen")
        os.makedirs(out_dir, exist_ok=True)
        base_name = f"Bewerbung_{company}_{datetime.now().strftime('%Y-%m-%d')}"
        out_path  = os.path.join(out_dir, base_name + ".pdf")
        counter   = 1
        while os.path.isfile(out_path):
            out_path = os.path.join(out_dir, f"{base_name}_{counter}.pdf")
            counter += 1

        self._bew_out_dir          = out_dir
        self._current_bew_job_id   = job["id"]
        self._current_bew_out_path = out_path
        self._bew_creating = True
        self._btn_create.setEnabled(False)

        replacements = {
            "{{JOBTITEL}}":  title,
            "{{ORT_DATUM}}": f"Flensburg den {today}",
        }
        pptx_src = self._pptx_path

        def _work():
            try:
                from pptx import Presentation
            except ImportError:
                self._bew_progress_sig.emit(0, "Error", "python-pptx not installed"); return

            self._bew_progress_sig.emit(15, "Loading template…", os.path.basename(pptx_src))
            prs = Presentation(pptx_src)

            self._bew_progress_sig.emit(35, "Replacing placeholders…", f"Job: {title[:55]}")
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        full = "".join(r.text for r in para.runs)
                        replaced = full
                        for ph, val in replacements.items():
                            replaced = replaced.replace(ph, val)
                        if replaced != full and para.runs:
                            para.runs[0].text = replaced
                            for r in para.runs[1:]:
                                r.text = ""

            soffice = self._find_soffice()
            if not soffice:
                self._bew_progress_sig.emit(0, "Error", "LibreOffice not found"); return

            with tempfile.TemporaryDirectory() as tmp:
                tmp_pptx = os.path.join(tmp, "bew_tmp.pptx")
                self._bew_progress_sig.emit(55, "Saving modified template…", "Writing PPTX…")
                prs.save(tmp_pptx)

                self._bew_progress_sig.emit(-1, "Converting to PDF…",
                    "LibreOffice is rendering your application — this takes a few seconds…")
                result = subprocess.run(
                    [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp, tmp_pptx],
                    capture_output=True, timeout=90,
                )
                tmp_pdf = os.path.join(tmp, "bew_tmp.pdf")
                if result.returncode != 0 or not os.path.isfile(tmp_pdf):
                    err = result.stderr.decode(errors="ignore")[:120]
                    self._bew_progress_sig.emit(0, "Conversion failed", err); return

                self._bew_progress_sig.emit(90, "Saving PDF…", out_path)
                shutil.copy2(tmp_pdf, out_path)

            self._bew_progress_sig.emit(100, "✔  Application ready!", os.path.basename(out_path))
            self._bew_pdf_ready.emit(out_path)

        threading.Thread(target=_work, daemon=True).start()

    # ── job profile HTML ──────────────────────────────────────────────────────

    def _build_desc_html(self, job: dict, loading: bool = False) -> str:
        c = P
        reason = html.escape(job.get("relevance_reason") or "")
        reason_block = (
            f'<div style="margin:0 0 18px 0;padding:10px 14px;background:{c["indigo_bg"]};'
            f'border-radius:8px;border-left:3px solid {c["indigo"]};">'
            f'<div style="color:{c["indigo"]};font-size:9px;font-weight:700;'
            f'letter-spacing:1px;margin-bottom:4px;">RELEVANCE REASON</div>'
            f'<div style="color:{c["text2"]};font-size:12px;">{reason}</div></div>'
        ) if reason else ""

        if loading:
            body = (
                f'<div style="color:{c["text3"]};font-size:13px;font-style:italic;padding:8px 0;">'
                f'● Loading description…</div>'
            )
        else:
            desc_plain = strip_html(job.get("description") or "")
            if len(desc_plain.strip()) > 30:
                body = _format_desc_html(desc_plain, c, job.get("title", ""))
            else:
                url = html.escape(job.get("url") or "")
                body = (
                    f'<div style="color:{c["text3"]};font-size:13px;">'
                    f'No description available.<br><br>'
                    f'<a href="{url}" style="color:{c["indigo"]};">'
                    f'Open job posting →</a></div>'
                )

        return (
            f'<html><body style="margin:16px 20px;font-family:\'Segoe UI\',sans-serif;'
            f'background:{c["card"]};color:{c["text2"]};">'
            f'{reason_block}{body}</body></html>'
        )

    # ── job list ──────────────────────────────────────────────────────────────

    def _funnel_html(self, counts: dict) -> str:
        c = P
        f = f"font-family:{MONO_FONT};font-size:12px;"
        return (
            f'<span style="{f}color:{c["text3"]};">PIPELINE  </span>'
            f'<span style="{f}color:{c["text2"]};">NEW</span> '
            f'<b style="{f}color:{c["text2"]};">{counts.get("new", 0)}</b>'
            f'<span style="{f}color:{c["text3"]};"> · </span>'
            f'<span style="{f}color:{c["indigo"]};">APPLIED</span> '
            f'<b style="{f}color:{c["indigo"]};">{counts.get("applied", 0)}</b>'
            f'<span style="{f}color:{c["text3"]};"> · </span>'
            f'<span style="{f}color:{c["amber"]};">INTERVIEW</span> '
            f'<b style="{f}color:{c["amber"]};">{counts.get("interview", 0)}</b>'
            f'<span style="{f}color:{c["text3"]};"> · </span>'
            f'<span style="{f}color:{c["green"]};">OFFER</span> '
            f'<b style="{f}color:{c["green"]};">{counts.get("offer", 0)}</b>'
        )

    def _load_jobs(self):
        cat = None if self._filter_cat in ("All", "All Categories") else self._filter_cat.lower()
        self._jobs = db.get_jobs(
            min_score=self._score_slider.value(),
            category=cat,
            show_dismissed=self._show_dismissed,
            view=self._view,
            search_text=self._search_input.text().strip(),
            sort=self._sort_combo.currentData(),
            ai_only=self._chk_ai_only.isChecked(),
            new_only=self._quick_new_only,
            status_filter=self._quick_status,
        )
        self._lbl_funnel.setText(self._funnel_html(db.get_pipeline_counts()))
        self._render_list()
        self._total_jobs_cache = len(db.get_jobs(show_dismissed=True))
        col = P['amber'] if self._total_jobs_cache > 0 else P['text3']
        self._st_total.setText(str(self._total_jobs_cache))
        self._st_total.setStyleSheet(
            f"color: {col}; font-size: 18px; font-weight: 700; "
            f"font-family: '{MONO_FONT}'; background: transparent;"
        )
        saved_last = search_engine.search_status.get("saved", 0)
        self._lbl_idle_db.setText(
            f"DB  {self._total_jobs_cache} jobs" + (f"  ·  {saved_last} new last run" if saved_last else "")
        )

    def _render_list(self):
        while self._list_layout.count() > 1:
            item = self._list_layout.takeAt(0)
            if item.widget(): item.widget().deleteLater()
        self._cards.clear()
        n = len(self._jobs)
        self._lbl_count.setText(f"{n} jobs")
        nd = len([j for j in self._jobs if not j.get("dismissed")])
        self._btn_dismiss_all.setText(f"✕ Dismiss ({nd})" if nd else "✕ Dismiss")
        self._lbl_count.setStyleSheet(
            f"color: {P['green'] if n > 0 else P['red']}; "
            f"font-size: 10px; font-family: '{MONO_FONT}';"
        )
        for job in self._jobs:
            card = JobCard(job)
            if self._selected and self._selected["id"] == job["id"]:
                card.set_selected(True)
            card.clicked.connect(self._select)
            self._list_layout.insertWidget(self._list_layout.count() - 1, card)
            self._cards[job["id"]] = card
        # auto-select first if none selected or selection gone
        ids = {j["id"] for j in self._jobs}
        if self._jobs and (self._selected is None or self._selected["id"] not in ids):
            self._select(self._jobs[0])

    # ── selection & detail ────────────────────────────────────────────────────

    def _select(self, job: dict):
        # flush pending notes before switching jobs
        if self._selected and self._selected["id"] != job["id"]:
            if self._notes_timer.isActive():
                self._notes_timer.stop()
                self._save_notes()
        if self._selected and self._selected["id"] in self._cards:
            self._cards[self._selected["id"]].set_selected(False)
        self._selected = job
        if job["id"] in self._cards:
            self._cards[job["id"]].set_selected(True)
        self._show_detail(job)
        self._maybe_fetch_desc(job)
        self._update_bew_status()

    def _maybe_fetch_desc(self, job: dict):
        """For BA jobs with short/empty descriptions, fetch full text in background."""
        if job.get("source") != "ba":
            return
        if len((job.get("description") or "").strip()) > 200:
            return
        refnr = (job.get("url") or "").rstrip("/").split("/")[-1]
        if not refnr or len(refnr) < 5:
            return
        job_id = job["id"]
        self._txt_desc.setHtml(self._build_desc_html(job, loading=True))
        _log_activity(f"Fetching description for {refnr}", "info")

        def _run():
            desc = ""
            try:
                from backend.services.ba_fetch import fetch_description
                desc = asyncio.run(fetch_description(refnr))
                _log_activity(f"BA desc fetched: {len(desc)} chars for {refnr}", "info")
            except Exception as e:
                _log_activity(f"BA desc error ({refnr}): {e}", "info")
            self._ba_desc_ready.emit(job_id, desc)  # thread-safe signal

        threading.Thread(target=_run, daemon=True).start()

    def _ba_desc_done(self, job_id: int, desc: str):
        if desc:
            db.update_job_description(job_id, desc)
        if self._selected and self._selected["id"] == job_id:
            if desc:
                self._selected["description"] = desc
            self._txt_desc.setHtml(self._build_desc_html(self._selected))

    def _update_detail_header(self, job: dict):
        if not job.get("viewed"):
            db.mark_viewed(job["id"])
            job["viewed"] = True
            if job["id"] in self._cards:
                self._cards[job["id"]].job["viewed"] = True
                self._cards[job["id"]]._apply_style()
        self._lbl_title.setText(job["title"])
        parts = [p for p in [job.get("company"), job.get("location")] if p]
        if job.get("posted_at"):
            parts.append(str(job["posted_at"])[:10])
        self._lbl_meta.setText("  ·  ".join(parts))

        score = job.get("relevance_score")
        ai_scored = (job.get("relevance_reason") or "").startswith("[AI]")
        fg, bg = score_col(score)
        score_text = (f"✦ {score}" if ai_scored else str(score)) if score is not None else "—"
        self._lbl_score.setText(score_text)
        self._lbl_score.setStyleSheet(
            f"background: {bg}; color: {fg}; border-radius: 8px; "
            f"padding: 5px 14px; font-size: 16px; font-weight: 700;"
        )
        self._lbl_reason.setText(job.get("relevance_reason") or "")

        while self._tag_row.count():
            item = self._tag_row.takeAt(0)
            if item.widget(): item.widget().deleteLater()
        cat = job.get("category", "")
        if cat in CAT: self._tag_row.addWidget(_pill(cat.upper(), *CAT[cat]))
        src = job.get("source", "")
        if src: self._tag_row.addWidget(_pill(src.upper(), P['text3'], P['border']))
        self._tag_row.addStretch()

        self._txt_desc.setHtml(self._build_desc_html(job))
        self._txt_desc.verticalScrollBar().setValue(0)

    def _update_detail_state(self, job: dict):
        self._tabs.setCurrentIndex(self._active_tab)

        self._refresh_save_btn(job.get("saved", False))
        self._btn_dismiss.setText("Restore" if job.get("dismissed", False) else "Dismiss")
        self._update_status_pipeline(job.get("status", "new"))

        self._txt_notes.blockSignals(True)
        self._txt_notes.setPlainText(db.get_notes(job["id"]) or "")
        self._txt_notes.blockSignals(False)
        self._lbl_notes_saved.setText("")

    def _show_detail(self, job: dict):
        self._update_detail_header(job)
        self._update_detail_state(job)

    def _clear_detail(self):
        self._lbl_title.setText("Select a job from the list")
        self._lbl_meta.setText(""); self._lbl_reason.setText("")
        self._lbl_score.setText("—")
        self._lbl_score.setStyleSheet(
            f"background: {P['border']}; color: {P['text3']}; border-radius: 8px; "
            f"padding: 5px 14px; font-size: 16px; font-weight: 700;"
        )
        while self._tag_row.count():
            item = self._tag_row.takeAt(0)
            if item.widget(): item.widget().deleteLater()
        self._txt_desc.clear()
        self._update_status_pipeline("new")
        self._txt_notes.blockSignals(True); self._txt_notes.clear(); self._txt_notes.blockSignals(False)

    # ── actions ───────────────────────────────────────────────────────────────

    def _open_job(self):
        if self._selected: webbrowser.open(self._selected["url"])

    def _refresh_save_btn(self, saved: bool):
        self._btn_save.setText("● Saved" if saved else "Save")
        self._btn_save.setStyleSheet(
            f"QPushButton {{ background: {P['green_bg'] if saved else P['card2']}; "
            f"color: {P['green'] if saved else P['text']}; "
            f"border-radius: 8px; font-size: 12px; font-weight: 600; padding: 6px 14px; }}"
            f"QPushButton:hover {{ background: {'#0a2a18' if saved else P['card3']}; }}"
        )

    def _toggle_save(self):
        if not self._selected: return
        saved = db.toggle_save(self._selected["id"])
        self._selected["saved"] = saved
        self._refresh_save_btn(saved)
        jid = self._selected["id"]
        if jid in self._cards:
            self._cards[jid].job["saved"] = saved
            self._cards[jid]._rebuild_tags()
            self._cards[jid]._apply_style()

    def _toggle_dismiss(self):
        if not self._selected: return
        self._last_dismissed_id = self._selected["id"]
        self._btn_undo_dis.setEnabled(True)
        try:
            cur_idx = next(i for i, j in enumerate(self._jobs) if j["id"] == self._selected["id"])
        except StopIteration:
            cur_idx = 0
        db.toggle_dismiss(self._selected["id"])
        self._selected = None
        self._load_jobs()
        if self._jobs:
            self._select(self._jobs[min(cur_idx, len(self._jobs) - 1)])
        else:
            self._clear_detail()

    def _undo_dismiss(self):
        if self._last_dismissed_id is None: return
        db.undo_dismiss(self._last_dismissed_id)
        self._last_dismissed_id = None
        self._btn_undo_dis.setEnabled(False)
        self._load_jobs()

    def _rescore_current(self):
        if not self._selected: return
        if ai_score_engine.ai_score_status["running"]: return
        s = db.get_settings()
        model   = s["prefs"].get("ollama_model", "qwen2.5:14b")
        profile = _profile_text(s["profile"])
        prompt  = s["prefs"].get("ai_score_prompt", "")
        job     = self._selected
        self._btn_rescore.setText("✦ …")
        self._btn_rescore.setEnabled(False)

        def _run():
            from backend.services.ai_scorer import score_job_ai
            result = score_job_ai(
                job.get("title",""), job.get("company",""),
                job.get("description",""), profile, model, prompt,
            )
            if result["ok"]:
                db.update_job_ai_score(job["id"], result["score"], result["reason"])
                self._ai_job_ready.emit(job["id"], result["score"], result["reason"])
            QTimer.singleShot(0, lambda: (
                self._btn_rescore.setText("✦ Re-score"),
                self._btn_rescore.setEnabled(True),
            ))

        threading.Thread(target=_run, daemon=True).start()

    def _update_status_pipeline(self, status: str):
        self._btn_applied_pipe.setStyleSheet(_pipe_qss(P['indigo'], status == "applied"))
        self._btn_interview.setStyleSheet(   _pipe_qss(P['amber'],  status == "interview"))
        self._btn_offer.setStyleSheet(       _pipe_qss(P['green'],  status == "offer"))
        self._btn_rejected.setStyleSheet(    _pipe_qss(P['red'],    status == "rejected"))

    def _set_status_pipeline(self, status: str):
        if not self._selected: return
        new_status = db.set_status(self._selected["id"], status)
        self._selected["status"] = new_status
        self._update_status_pipeline(new_status)
        jid = self._selected["id"]
        if jid in self._cards:
            self._cards[jid].job["status"] = new_status
            self._cards[jid]._rebuild_tags()
            self._cards[jid]._apply_style()

    def _on_notes_changed(self):
        self._notes_timer.start(1000)
        self._lbl_notes_saved.setText("editing…")

    def _save_notes(self):
        if not self._selected: return
        text = self._txt_notes.toPlainText()
        db.save_notes(self._selected["id"], text)
        self._selected["notes"] = text
        self._lbl_notes_saved.setText("Saved")
        QTimer.singleShot(2000, lambda: self._lbl_notes_saved.setText(""))

    # ── search ────────────────────────────────────────────────────────────────

    def _start_search(self):
        if search_engine.search_status["running"]: return
        self._search_start = time.time()
        self._btn_search.setText("Searching…"); self._btn_search.setEnabled(False)
        self._btn_cancel.setEnabled(True)
        self._progress.setRange(0, 100); self._progress.setValue(0)
        self._progress.setFormat("  ● Initialising…")
        search_engine.start_search(on_done=lambda: QTimer.singleShot(0, self._search_done))

    def _cancel_search(self):
        search_engine.cancel_search(); self._btn_cancel.setEnabled(False)

    # ── AI scoring ────────────────────────────────────────────────────────────

    def _on_ai_toggle(self):
        self._ai_auto_enabled = self._btn_ai_toggle.isChecked()
        self._apply_ai_toggle_style()

    def _apply_ai_toggle_style(self):
        on = self._ai_auto_enabled
        self._btn_ai_toggle.setText("AUTO  ON " if on else "AUTO  OFF")
        self._btn_ai_toggle.setStyleSheet(
            f"QPushButton {{ background: {P['purple_bg'] if on else P['card2']}; "
            f"color: {P['purple'] if on else P['text3']}; border-radius: 6px; "
            f"border: 1px solid {P['purple'] if on else P['border']}; "
            f"font-size: 11px; font-weight: 700; padding: 0 10px; }}"
            f"QPushButton:hover {{ background: {'#2a0e50' if on else P['card3']}; }}"
        )

    def _start_ai_scoring(self):
        if ai_score_engine.ai_score_status["running"]: return
        s    = db.get_settings()
        jobs = db.get_jobs_for_ai_scoring(min_score=int(s["prefs"].get("ai_min_score", 40)))
        if not jobs:
            self._lbl_ai_phase.setText("nothing to score")
            return
        model   = s["prefs"].get("ollama_model", "qwen2.5:14b")
        profile = _profile_text(s["profile"])
        prompt  = s["prefs"].get("ai_score_prompt", "")
        self._btn_ai_start.setEnabled(False)
        self._btn_ai_cancel.setEnabled(True)
        _log_activity(f"AI scoring started: {len(jobs)} jobs queued", "info")
        ai_score_engine.start_scoring(
            jobs, model, profile, prompt,
            on_job_done =lambda jid, sc, rs: self._ai_job_ready.emit(jid, sc, rs),
            on_complete =lambda: QTimer.singleShot(0, self._ai_scoring_done),
        )

    def _ai_job_done(self, job_id: int, score: int, reason: str):
        for j in self._jobs:
            if j["id"] == job_id:
                j["relevance_score"] = score; j["relevance_reason"] = reason; break
        if score < 50:
            db.force_dismiss(job_id)
            for j in self._jobs:
                if j["id"] == job_id:
                    j["dismissed"] = True; break
            if self._selected and self._selected["id"] == job_id:
                self._selected = None
            self._ai_render_timer.start(800)
            return
        if self._selected and self._selected["id"] == job_id:
            self._selected["relevance_score"] = score
            self._selected["relevance_reason"] = reason
            self._show_detail(self._selected)
        self._ai_render_timer.start(800)

    def _ai_scoring_done(self):
        self._btn_ai_start.setEnabled(True)
        self._btn_ai_cancel.setEnabled(False)
        scored = ai_score_engine.ai_score_status.get("scored", 0)
        _log_activity(f"AI scoring complete — {scored} jobs updated", "info")
        self._load_jobs()
        if scored > 0:
            self._tray.showMessage(
                "AutoApply", f"AI scoring complete — {scored} jobs rated!",
                QSystemTrayIcon.MessageIcon.Information, 4000,
            )

    def _slbl(self, lbl, text: str, color: str):
        lbl.setText(text)
        lbl.setStyleSheet(f"color: {color}; font-size: 11px; font-family: '{MONO_FONT}';")

    def _poll_ai(self):
        st = ai_score_engine.ai_score_status
        running = st["running"]
        total   = st["total"]
        done    = st["done"]
        phase   = st["phase"]
        eta_s   = st["eta_s"]
        current = st["current"]

        if running and total > 0:
            pct = int(done / total * 100)
            self._ai_progress.setValue(pct)
            self._ai_progress.setFormat(f"  ✦ {done}/{total}  ·  {pct}%")
            self._btn_ai_cancel.setEnabled(True)
            self._btn_ai_start.setEnabled(False)
            eta_m, eta_s2 = divmod(eta_s, 60)
            self._lbl_ai_eta.setText(f"ETA {eta_m}:{eta_s2:02d}")
            self._slbl(self._lbl_ai_phase, f"✦ {current}" if current else "✦ scoring…", P['purple'])
        elif phase == "done":
            scored = st.get("scored", 0)
            self._ai_progress.setValue(100)
            self._ai_progress.setFormat(f"  ✦ Done · {scored} updated")
            self._lbl_ai_eta.setText("")
            self._slbl(self._lbl_ai_phase, "done", P['green'])
        elif phase == "cancelling":
            self._ai_progress.setFormat("  cancelling…")
            self._slbl(self._lbl_ai_phase, "CANCELLING…", P['red'])
            self._lbl_ai_eta.setText("")
        elif phase == "cancelled":
            self._ai_progress.setFormat("  cancelled")
            self._slbl(self._lbl_ai_phase, "cancelled", P['red'])
            self._lbl_ai_eta.setText("")
        else:
            # throttle idle DB query to every ~5s (poll runs every 400ms → every 12 ticks)
            self._ai_idle_tick += 1
            if self._ai_idle_tick % 12 == 1:
                self._ai_unscored_cache = len(db.get_jobs_for_ai_scoring(40))
            self._ai_progress.setValue(0)
            self._ai_progress.setFormat(f"  IDLE · {self._ai_unscored_cache} queued" if self._ai_unscored_cache else "  IDLE")
            self._lbl_ai_eta.setText("")
            self._slbl(self._lbl_ai_phase, "idle", P['text3'])

    def _search_done(self):
        self._btn_search.setText("Start Search"); self._btn_search.setEnabled(True)
        self._btn_cancel.setEnabled(False)
        saved = search_engine.search_status.get("saved", 0)
        self._progress.setValue(100)
        self._progress.setFormat(f"  ✓  {saved} new jobs saved")
        self._lbl_search_meta.setText("")
        QTimer.singleShot(3000, lambda: (
            self._progress.setValue(0),
            self._progress.setFormat("  READY"),
        ))
        self._run_auto_dismiss()
        n_empty = db.dismiss_empty_description_jobs()
        if n_empty > 0:
            _log_activity(f"Auto-dismissed {n_empty} jobs without description", "db")
        if saved > 0:
            self._chip_new.setChecked(True)
            self._quick_new_only = True
            self._score_slider.blockSignals(True)
            self._score_slider.setValue(40)
            self._score_slider.blockSignals(False)
            self._lbl_min_score.setText("MIN SCORE  40")
        self._load_jobs(); self._update_last_search_label()
        if saved > 0:
            self._tray.showMessage("AutoApply", f"Search complete — {saved} new jobs found!",
                                   QSystemTrayIcon.MessageIcon.Information, 4000)
        if self._ai_auto_enabled:
            QTimer.singleShot(2000, self._start_ai_scoring)

    def _poll_search(self):
        st      = search_engine.search_status
        running = st.get("running", False)
        phase   = st.get("phase", "IDLE")
        fetched = st.get("fetched", 0)
        filtered = st.get("filtered", 0)
        saved   = st.get("saved", 0)

        if running:
            pct = _search_progress_pct(phase)
            self._progress.setValue(pct)
            self._progress.setFormat(f"  {pct}%  ·  {fetched} jobs fetched")
            elapsed = int(time.time() - getattr(self, "_search_start", time.time()))
            mins, secs = divmod(elapsed, 60)
            rate = fetched / max(elapsed, 1)
            self._lbl_search_meta.setText(f"{mins}:{secs:02d}  ·  {rate:.1f} j/s")
            self._btn_search.setEnabled(False); self._btn_search.setText("Searching…")
            self._btn_cancel.setEnabled(True)

        self._stats_frame.setVisible(running)
        self._lbl_idle_db.setVisible(not running)

        self._slbl(self._lbl_phase, ("● " if running else "") + phase.upper(),
                   P['green'] if running else P['text3'])

        def _val_style(col):
            return f"color: {col}; font-size: 20px; font-weight: 700; font-family: '{MONO_FONT}'; background: transparent;"

        self._st_fetched.setText(str(fetched) if fetched else "—")
        self._st_fetched.setStyleSheet(_val_style(
            P['green'] if fetched > 0 else (P['red'] if st.get('done') else P['text3'])
        ))
        self._st_filtered.setText(str(filtered) if filtered else "—")
        self._st_filtered.setStyleSheet(_val_style(P['amber'] if filtered > 0 else P['text3']))
        self._st_saved.setText(str(saved) if saved else "—")
        self._st_saved.setStyleSheet(_val_style(P['indigo'] if saved > 0 else P['text3']))
        self._update_topbar(running)

    def _update_topbar(self, running: bool):
        self._top_time.setText(time.strftime("%H:%M:%S"))
        if running:
            self._top_status.setText("● SEARCHING")
            self._top_status.setStyleSheet(f"color: {P['indigo']}; font-size: 10px; font-family: '{MONO_FONT}';")
            self._top_dot.setStyleSheet(f"color: {P['indigo']}; font-size: 10px; background: transparent;")
        else:
            self._top_status.setText("○ IDLE")
            self._top_status.setStyleSheet(f"color: {P['text3']}; font-size: 10px; font-family: '{MONO_FONT}';")
            self._top_dot.setStyleSheet(f"color: {P['green']}; font-size: 10px; background: transparent;")
        self._top_jobs_tick += 1
        if self._top_jobs_tick % 5 == 1:
            self._top_jobs.setText(f"DB  {self._total_jobs_cache} jobs")

    # ── filter / view ─────────────────────────────────────────────────────────

    def _on_score_slider(self, value: int):
        self._lbl_min_score.setText(f"MIN SCORE  {value}")
        self._load_jobs()

    def _on_view_change(self, view: str):
        self._view = view.lower(); self._load_jobs()

    def _on_filter(self, *_):
        self._filter_cat     = self._cat_combo.currentText()
        self._show_dismissed = self._chk_dismissed.isChecked()
        self._load_jobs()

    def _on_chip_new(self):
        self._quick_new_only = self._chip_new.isChecked()
        self._load_jobs()

    def _on_chip_interview(self):
        self._quick_status = "interview" if self._chip_interview.isChecked() else ""
        self._load_jobs()

    def _run_auto_dismiss(self):
        days = int(db.get_settings()["prefs"].get("auto_dismiss_days", 0))
        if days <= 0:
            return
        n = db.auto_dismiss_old_jobs(days)
        if n > 0:
            _log_activity(f"Auto-dismissed {n} jobs older than {days} days", "db")

    def _dismiss_visible(self):
        to_dismiss = [j for j in self._jobs if not j.get("dismissed")]
        if not to_dismiss:
            return
        reply = QMessageBox.question(
            self, "Dismiss visible jobs",
            f"Dismiss all {len(to_dismiss)} visible jobs?",
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.Cancel,
        )
        if reply != QMessageBox.StandardButton.Yes:
            return
        db.bulk_dismiss([j["id"] for j in to_dismiss])
        _log_activity(f"Bulk dismissed {len(to_dismiss)} jobs", "db")
        self._load_jobs()

    def _export_csv(self):
        path, _ = QFileDialog.getSaveFileName(self, "Export Jobs", "jobs.csv", "CSV (*.csv)")
        if not path:
            return
        fields = ["title", "company", "location", "relevance_score", "status", "url", "relevance_reason"]
        with open(path, "w", newline="", encoding="utf-8-sig") as f:
            w = csv.DictWriter(f, fieldnames=fields, extrasaction="ignore")
            w.writeheader()
            w.writerows(self._jobs)
        _log_activity(f"Exported {len(self._jobs)} jobs to CSV", "db")

    def _clear_jobs(self):
        reply = QMessageBox.question(
            self, "Clear all jobs?",
            "This will permanently delete all jobs, applications and cover letters.\n\nContinue?",
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.Cancel,
        )
        if reply == QMessageBox.StandardButton.Yes:
            db.clear_all_jobs()
            self._selected = None; self._clear_detail(); self._load_jobs()
            _log_activity("Job list cleared by user", "db")

    # ── auto-search & tray ────────────────────────────────────────────────────

    def _auto_search_check(self):
        if search_engine.search_status["running"]: return
        s     = db.get_settings()
        hours = int(s["prefs"].get("auto_search_hours", 0))
        if hours <= 0: return
        last  = s["prefs"].get("last_search_ts", "")
        if not last:
            self._start_search(); return
        try:
            for fmt in ("%Y-%m-%d %H:%M:%S", "%Y-%m-%d %H:%M"):
                try:
                    last_dt = datetime.strptime(last, fmt); break
                except ValueError:
                    continue
            else:
                return
            if (datetime.utcnow() - last_dt).total_seconds() / 3600 >= hours:
                self._start_search()
        except Exception:
            pass

    def _update_last_search_label(self):
        s   = db.get_settings()
        ts  = s["prefs"].get("last_search_ts", "")
        new = s["prefs"].get("last_search_new", 0)
        self._lbl_last_search.setText(
            f"Last search: {ts}  ·  {new} new" if ts else "Last search: never"
        )

    # ── keyboard shortcuts ────────────────────────────────────────────────────

    def _next_unviewed(self):
        if not self._jobs: return
        start = 0
        if self._selected:
            try:
                start = next(i for i, j in enumerate(self._jobs) if j["id"] == self._selected["id"]) + 1
            except StopIteration:
                start = 0
        for i in range(start, len(self._jobs)):
            if not self._jobs[i].get("viewed"):
                self._select(self._jobs[i]); return
        for i in range(0, start):
            if not self._jobs[i].get("viewed"):
                self._select(self._jobs[i]); return
        if start < len(self._jobs):
            self._select(self._jobs[start])

    def _go_to_bewerbung(self):
        self._tabs.setCurrentIndex(1)
        self._create_application_pdf()

    def keyPressEvent(self, event):
        key = event.key()
        if   key == Qt.Key.Key_S: self._toggle_save()
        elif key == Qt.Key.Key_D: self._toggle_dismiss()
        elif key == Qt.Key.Key_A: self._set_status_pipeline("applied")
        elif key == Qt.Key.Key_B: self._go_to_bewerbung()
        elif key == Qt.Key.Key_N: self._tabs.setCurrentIndex(2)
        elif key == Qt.Key.Key_O: self._open_job()
        elif key == Qt.Key.Key_Tab: self._next_unviewed()
        elif key in (Qt.Key.Key_Down, Qt.Key.Key_Up):
            if not self._jobs: return
            if self._selected is None:
                idx = 0
            else:
                try:
                    cur = next(i for i, j in enumerate(self._jobs) if j["id"] == self._selected["id"])
                    idx = min(cur + 1, len(self._jobs) - 1) if key == Qt.Key.Key_Down else max(cur - 1, 0)
                except StopIteration:
                    idx = 0
            self._select(self._jobs[idx])
        else:
            super().keyPressEvent(event)

    # ── dialogs ───────────────────────────────────────────────────────────────

    def _open_stats(self):
        StatsDialog(self).exec()

    def _open_debug(self):
        if self._debug_win is None:
            self._debug_win = DebugWindow(self)
        self._debug_win.show(); self._debug_win.raise_(); self._debug_win.activateWindow()

    def _open_settings(self):
        SettingsDialog(self).exec()


# ── settings dialog ───────────────────────────────────────────────────────────

class StatsDialog(QDialog):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Statistics"); self.resize(520, 520)
        self.setStyleSheet(f"QDialog {{ background: {P['bg']}; }}")
        self._build()

    def _build(self):
        stats = db.get_stats()
        lay = QVBoxLayout(self); lay.setContentsMargins(28, 24, 28, 24); lay.setSpacing(18)
        lay.addWidget(_label("Statistics", size=16, color=P['text'], bold=True))
        lay.addWidget(_sep())

        def _bar_row(label: str, value: int, total: int, color: str):
            w = QWidget(); hl = QHBoxLayout(w); hl.setContentsMargins(0,0,0,0); hl.setSpacing(10)
            lbl = _mono(label, size=11, color=P['text2']); lbl.setFixedWidth(90)
            pct = value / max(total, 1)
            bar_bg = QFrame(); bar_bg.setFixedHeight(16); bar_bg.setStyleSheet(
                f"background:{P['card2']}; border-radius:4px;")
            bar_fill = QFrame(bar_bg); bar_fill.setFixedHeight(16)
            bar_fill.setFixedWidth(max(4, int(260 * pct)))
            bar_fill.setStyleSheet(f"background:{color}; border-radius:4px;")
            cnt = _mono(str(value), size=11, color=color); cnt.setFixedWidth(40)
            hl.addWidget(lbl); hl.addWidget(bar_bg, 1); hl.addWidget(cnt)
            return w

        total = max(stats["total"], 1)

        lay.addWidget(_label("Score distribution", size=12, color=P['text2'], bold=True))
        for band, col in [("0-25", P['red']), ("25-50", P['amber']), ("50-75", P['indigo']), ("75-100", P['green'])]:
            lay.addWidget(_bar_row(band, stats["score_bands"][band], total, col))

        lay.addWidget(_sep())
        lay.addWidget(_label("Categories", size=12, color=P['text2'], bold=True))
        cat_cols = {"it": CAT["it"][0], "wirtschaft": CAT["wirtschaft"][0], "unknown": P['text3']}
        for cat, col in cat_cols.items():
            lay.addWidget(_bar_row(cat.capitalize(), stats["categories"].get(cat, 0), total, col))

        lay.addWidget(_sep())
        lay.addWidget(_label("Application pipeline", size=12, color=P['text2'], bold=True))
        pipe_cols = [("new", P['text3']), ("applied", P['indigo']), ("interview", P['amber']),
                     ("offer", P['green']), ("rejected", P['red'])]
        pipe_total = max(sum(stats["pipeline"].values()), 1)
        for st, col in pipe_cols:
            lay.addWidget(_bar_row(st.capitalize(), stats["pipeline"].get(st, 0), pipe_total, col))

        lay.addWidget(_sep())
        summary_row = QHBoxLayout()
        for lbl, val, col in [
            ("Total jobs", stats["total"], P['text']),
            ("AI scored",  stats["ai_scored"], P['purple']),
            ("Viewed",     stats["viewed"], P['indigo']),
        ]:
            cell = QVBoxLayout()
            cell.addWidget(_label(str(val), size=22, color=col, bold=True))
            cell.addWidget(_label(lbl, size=10, color=P['text3']))
            summary_row.addLayout(cell)
        lay.addLayout(summary_row)
        lay.addStretch()

        close = _btn("Close", P['card2'], P['card3']); close.clicked.connect(self.accept)
        lay.addWidget(close)


class SettingsDialog(QDialog):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Settings"); self.resize(680, 720)
        self.setStyleSheet(f"QDialog {{ background: {P['bg']}; }}")
        self._s = db.get_settings()
        self._build()

    def _build(self):
        lay = QVBoxLayout(self); lay.setContentsMargins(24, 24, 24, 24); lay.setSpacing(16)
        lay.addWidget(_label("Settings", size=16, color=P['text'], bold=True)); lay.addWidget(_sep())
        tabs = QTabWidget(); tabs.setDocumentMode(True); lay.addWidget(tabs, 1)

        # Profile
        p_w = QWidget(); p_w.setStyleSheet(f"background: {P['card']};")
        p_l = QFormLayout(p_w); p_l.setContentsMargins(20, 20, 20, 20); p_l.setSpacing(12)
        p_l.setLabelAlignment(Qt.AlignmentFlag.AlignRight)
        prof = self._s["profile"]
        self._f = {k: QLineEdit(prof.get(k, "")) for k in
                   ["name","email","phone","address","degree","background","motivation"]}
        for lbl, key in [("Name","name"),("Email","email"),("Phone","phone"),("Address","address"),
                          ("Degree","degree"),("Background","background"),("Motivation","motivation")]:
            p_l.addRow(_label(lbl, size=11, color=P['text2']), self._f[key])
        tabs.addTab(p_w, "Profile")

        # AI / Preferences
        ai_w = QWidget(); ai_w.setStyleSheet(f"background: {P['card']};")
        ai_l = QVBoxLayout(ai_w); ai_l.setContentsMargins(20, 20, 20, 20); ai_l.setSpacing(12)
        prefs = self._s["prefs"]

        def _pref_row(title, widget, hint=None, stretch=False):
            ai_l.addWidget(_label(title, size=11, color=P['text2'], bold=True))
            if hint:
                ai_l.addWidget(_label(hint, size=10, color=P['text3']))
            ai_l.addWidget(widget, 1 if stretch else 0)
            ai_l.addWidget(_sep())

        self._f_extra = QTextEdit(); self._f_extra.setPlainText(prefs.get("extra_prompt", ""))
        self._f_extra.setFixedHeight(90)
        _pref_row("Extra instructions for AI", self._f_extra)

        from backend.services.ai_scorer import _DEFAULT_CUSTOM
        self._f_ai_score_prompt = QTextEdit()
        self._f_ai_score_prompt.setPlainText(prefs.get("ai_score_prompt", "") or _DEFAULT_CUSTOM)
        self._f_ai_score_prompt.setFixedHeight(120)
        _pref_row("AI Scorer Prompt", self._f_ai_score_prompt,
                  hint="Zusätzliche Anweisungen für die KI beim Bewerten von Jobs (optional).")

        self._f_model = QLineEdit(prefs.get("ollama_model", "qwen2.5:14b"))
        _pref_row("Ollama model", self._f_model)

        _AUTO_OPTS = ["Off","Every 1h","Every 2h","Every 4h","Every 8h","Every 12h","Every 24h"]
        _AUTO_VALS = [0, 1, 2, 4, 8, 12, 24]
        self._f_auto = QComboBox(); self._f_auto.addItems(_AUTO_OPTS)
        try: self._f_auto.setCurrentIndex(_AUTO_VALS.index(int(prefs.get("auto_search_hours", 0))))
        except ValueError: self._f_auto.setCurrentIndex(0)
        _pref_row("Auto-Search interval", self._f_auto)

        self._f_ai_min_score = QLineEdit(str(prefs.get("ai_min_score", 40)))
        _pref_row("AI Scoring — Mindest-Regel-Score", self._f_ai_min_score,
                  hint="Nur Jobs mit Regel-Score ≥ X werden von der KI bewertet.")

        _DISMISS_OPTS = ["Off", "7 days", "14 days", "30 days", "60 days", "90 days"]
        _DISMISS_VALS = [0, 7, 14, 30, 60, 90]
        self._f_dismiss = QComboBox(); self._f_dismiss.addItems(_DISMISS_OPTS)
        try: self._f_dismiss.setCurrentIndex(_DISMISS_VALS.index(int(prefs.get("auto_dismiss_days", 0))))
        except ValueError: self._f_dismiss.setCurrentIndex(0)
        _pref_row("Auto-dismiss jobs older than", self._f_dismiss,
                  hint="Jobs mit Status 'new' älter als X Tage werden automatisch ausgeblendet.")

        tabs.addTab(ai_w, "AI / Preferences")

        btn_row = QHBoxLayout(); btn_row.addStretch()
        cancel = _btn("Cancel", P['card2'], P['card3']); cancel.clicked.connect(self.reject)
        save   = _btn("Save Settings", P['indigo'], P['indigo_d']); save.clicked.connect(self._save)
        btn_row.addWidget(cancel); btn_row.addSpacing(8); btn_row.addWidget(save)
        lay.addLayout(btn_row)

    def _save(self):
        _AUTO_VALS    = [0, 1, 2, 4, 8, 12, 24]
        _DISMISS_VALS = [0, 7, 14, 30, 60, 90]
        profile = {k: self._f[k].text() for k in self._f}
        prefs = {
            **self._s["prefs"],
            "extra_prompt":      self._f_extra.toPlainText(),
            "ai_score_prompt":   self._f_ai_score_prompt.toPlainText(),
            "ollama_model":      self._f_model.text().strip(),
            "auto_search_hours": _AUTO_VALS[self._f_auto.currentIndex()],
            "auto_dismiss_days": _DISMISS_VALS[self._f_dismiss.currentIndex()],
            "ai_min_score":      max(0, min(100, int(self._f_ai_min_score.text() or "40") if self._f_ai_min_score.text().strip().lstrip('-').isdigit() else 40)),
        }
        db.save_settings(profile, prefs)
        import backend.services.llm_service as svc
        lines = [
            f"Bewerber: {profile.get('name','')}",
            f"Abschluss: {profile.get('degree','')}",
            f"Hintergrund: {profile.get('background','')}",
            f"Motivation: {profile.get('motivation','')}",
        ]
        ctx = "\n".join(l for l in lines if l.strip())
        svc._PROFILE_CONTEXT = ctx
        svc._SYSTEM_PROMPT = (
            f"Du bist ein Bewerbungsassistent. Du kennst das folgende Bewerber-Profil:\n\n{ctx}\n\n"
            f"Antworte immer auf Deutsch. Sei präzise und halte dich an das geforderte Format."
        )
        self.accept()


# ── debug window ──────────────────────────────────────────────────────────────

class DebugWindow(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent, Qt.WindowType.Window)
        self.setWindowTitle("AutoApply  ·  Debug Console")
        self.resize(1000, 740); self.setMinimumSize(800, 600)
        self.setStyleSheet(f"QWidget {{ background: {P['bg']}; }}")
        self._tick = 0; self._dot_state = True
        self._build()
        self._timer = QTimer(self); self._timer.timeout.connect(self._refresh); self._timer.start(1000)
        try:
            import psutil; nc = psutil.net_io_counters()
            self._prev_net = (nc.bytes_sent, nc.bytes_recv)
        except Exception:
            self._prev_net = (0, 0)
        self._prev_llm_count = 0
        self._refresh()

    def _build(self):
        root = QVBoxLayout(self); root.setContentsMargins(20, 16, 20, 16); root.setSpacing(12)

        hdr = QHBoxLayout()
        title = QLabel("DEBUG CONSOLE")
        title.setStyleSheet(f"color: {P['text']}; font-size: 13px; font-weight: 700; letter-spacing: 2px;")
        hdr.addWidget(title); hdr.addStretch()
        self._lbl_dot  = _label("●  LIVE", size=10, color=P['green'])
        self._lbl_time = _label("", size=10, color=P['text3'])
        hdr.addWidget(self._lbl_dot); hdr.addSpacing(14); hdr.addWidget(self._lbl_time)
        root.addLayout(hdr); root.addWidget(_sep())

        # status cards
        cards = QHBoxLayout(); cards.setSpacing(10)
        self._sc_rows = self._status_card(cards, "SEARCH ENGINE", P["indigo"])
        self._ac_rows = self._status_card(cards, "AI / LLM",      P["purple"])
        self._dc_rows = self._status_card(cards, "DATABASE",      P["amber"])
        self._nc_rows = self._status_card(cards, "NETWORK",       P["green"])
        root.addLayout(cards)

        # graphs
        graphs = QHBoxLayout(); graphs.setSpacing(10)
        self._g_cpu   = self._graph_card(graphs, "CPU USAGE", P['indigo'], "%",      100)
        self._g_ram   = self._graph_card(graphs, "RAM USAGE", P['purple'], "%",      100)
        self._g_net_s = self._graph_card(graphs, "NET SEND",  P['green'],  " KB/s", 1024)
        self._g_net_r = self._graph_card(graphs, "NET RECV",  P['amber'],  " KB/s", 1024)
        root.addLayout(graphs)

        # logs
        bottom = QHBoxLayout(); bottom.setSpacing(10)

        act = _card(bg=P['card']); act_l = QVBoxLayout(act)
        act_l.setContentsMargins(14, 12, 14, 12); act_l.setSpacing(6)
        ah = QHBoxLayout()
        ah.addWidget(_label("ACTIVITY LOG", size=11, color=P['text3'], bold=True)); ah.addStretch()
        self._lbl_events = _label("", size=11, color=P['text3']); ah.addWidget(self._lbl_events)
        act_l.addLayout(ah)
        self._txt_activity = QTextEdit(); self._txt_activity.setReadOnly(True)
        self._txt_activity.setStyleSheet(
            f"QTextEdit {{ background: {P['card2']}; border: 1px solid {P['border']}; "
            f"border-radius: 8px; color: {P['green']}; padding: 8px; "
            f"font-family: '{MONO_FONT}'; font-size: 10px; }}"
        )
        act_l.addWidget(self._txt_activity, 1); bottom.addWidget(act, 3)

        # AI decisions panel
        aid = _card(bg=P['card']); aid_l = QVBoxLayout(aid)
        aid_l.setContentsMargins(14, 12, 14, 12); aid_l.setSpacing(6)
        dh = QHBoxLayout()
        dh.addWidget(_label("AI DECISIONS", size=11, color=P['text3'], bold=True)); dh.addStretch()
        self._lbl_ai_dec_count = _label("0 scored", size=11, color=P['purple'])
        dh.addWidget(self._lbl_ai_dec_count); aid_l.addLayout(dh)
        self._txt_ai_dec = QTextEdit(); self._txt_ai_dec.setReadOnly(True)
        self._txt_ai_dec.setStyleSheet(
            f"QTextEdit {{ background: {P['card2']}; border: 1px solid {P['border']}; "
            f"border-radius: 8px; color: {P['text2']}; padding: 8px; "
            f"font-family: '{MONO_FONT}'; font-size: 11px; }}"
        )
        aid_l.addWidget(self._txt_ai_dec, 1); bottom.addWidget(aid, 3)

        llm = _card(bg=P['card']); llm_l = QVBoxLayout(llm)
        llm_l.setContentsMargins(14, 12, 14, 12); llm_l.setSpacing(6)
        lh = QHBoxLayout()
        lh.addWidget(_label("LLM CALLS", size=11, color=P['text3'], bold=True)); lh.addStretch()
        self._lbl_llm_count = _label("0 calls", size=11, color=P['purple']); lh.addWidget(self._lbl_llm_count)
        llm_l.addLayout(lh)
        self._txt_llm = QTextEdit(); self._txt_llm.setReadOnly(True)
        self._txt_llm.setStyleSheet(
            f"QTextEdit {{ background: {P['card2']}; border: 1px solid {P['border']}; "
            f"border-radius: 8px; color: {P['text2']}; padding: 8px; "
            f"font-family: '{MONO_FONT}'; font-size: 10px; }}"
        )
        unload = _btn("Unload Model", P['red_bg'], "#400c0c", height=26, font_size=10)
        unload.clicked.connect(lambda: (
            threading.Thread(
                target=llm_service.unload_model,
                args=(db.get_settings()["prefs"].get("ollama_model", "qwen2.5:14b"),),
                daemon=True,
            ).start(),
            _log_activity("LLM model unloaded", "info"),
        ))
        llm_l.addWidget(self._txt_llm, 1); llm_l.addWidget(unload, 0, Qt.AlignmentFlag.AlignRight)
        bottom.addWidget(llm, 2)
        root.addLayout(bottom, 1)

    def _status_card(self, parent_lay: QHBoxLayout, title: str, accent: str) -> list:
        frame = _card(bg=P['card2']); lay = QVBoxLayout(frame)
        lay.setContentsMargins(14, 12, 14, 12); lay.setSpacing(5)
        hdr = QHBoxLayout()
        dot = QLabel("●"); dot.setStyleSheet(f"color: {accent}; font-size: 11px;")
        hdr.addWidget(dot); hdr.addWidget(_label(title, size=11, color=P['text3'], bold=True)); hdr.addStretch()
        lay.addLayout(hdr)
        rows = [_label("—", size=12, color=P['text2']) for _ in range(3)]
        for r in rows: lay.addWidget(r)
        lay.addStretch()
        parent_lay.addWidget(frame, 1)
        return rows

    def _graph_card(self, parent_lay: QHBoxLayout, title: str, color: str, unit: str, max_val=100) -> MiniGraph:
        frame = _card(bg=P['card2']); lay = QVBoxLayout(frame)
        lay.setContentsMargins(10, 8, 10, 8); lay.setSpacing(4)
        lay.addWidget(_label(title, size=12, color=P['text3'], bold=True))
        g = MiniGraph(color, title, unit, max_val); lay.addWidget(g, 1)
        parent_lay.addWidget(frame, 1)
        return g

    def _refresh(self):
        import psutil
        self._tick += 1
        self._lbl_time.setText(time.strftime("%Y-%m-%d  %H:%M:%S"))
        self._dot_state = not self._dot_state
        self._lbl_dot.setStyleSheet(
            f"color: {P['green'] if self._dot_state else P['text3']}; font-size: 10px;"
        )

        # search card
        st = search_engine.search_status
        running = st.get("running", False); phase = st.get("phase", "Ready")
        r = self._sc_rows
        r[0].setText(f"Status:   {'● RUNNING' if running else '○ idle'}")
        r[0].setStyleSheet(f"color: {P['green'] if running else P['text3']}; font-size: 12px;")
        r[1].setText(f"Phase:    {phase}"); r[1].setStyleSheet(f"color: {P['text2']}; font-size: 12px;")
        r[2].setText(f"Fetched: {st.get('fetched',0)}   New: {st.get('saved',0)}")
        r[2].setStyleSheet(f"color: {P['text3']}; font-size: 12px;")

        # AI / LLM merged card
        ast = ai_score_engine.ai_score_status
        log = llm_service.get_log(); last = log[-1] if log else None
        a = self._ac_rows
        a[0].setText(f"AI: {'● RUNNING' if ast['running'] else '○ ' + ast['phase']}  ·  LLM: {len(log)} calls")
        a[0].setStyleSheet(f"color: {P['purple'] if ast['running'] else P['text3']}; font-size: 12px;")
        eta = ast['eta_s']
        eta_str = f"ETA {eta//60}:{eta%60:02d}" if eta > 0 else f"avg {ast.get('avg_s', 0):.1f}s"
        a[1].setText(f"Scored: {ast.get('scored', 0)}/{ast['total']}  ·  {eta_str}")
        a[1].setStyleSheet(f"color: {P['text2']}; font-size: 12px;")
        last_txt = f"{last['purpose'][:22]}  ({last['duration_s']}s)" if last else "no calls yet"
        a[2].setText(f"Last: {last_txt}")
        a[2].setStyleSheet(f"color: {P['text3']}; font-size: 12px;")

        # db card — fetch once, derive counts from the result
        all_jobs = db.get_jobs(show_dismissed=True)
        total   = len(all_jobs)
        applied = sum(1 for j in all_jobs if j.get("status") == "applied")
        saved   = sum(1 for j in all_jobs if j.get("saved"))
        d = self._dc_rows
        d[0].setText(f"Total:    {total} jobs");  d[0].setStyleSheet(f"color: {P['amber']}; font-size: 12px;")
        d[1].setText(f"Saved:    {saved}");        d[1].setStyleSheet(f"color: {P['text2']}; font-size: 12px;")
        d[2].setText(f"Applied:  {applied}");      d[2].setStyleSheet(f"color: {P['text3']}; font-size: 12px;")

        # network card
        try:
            nc = psutil.net_io_counters()
            sent_kb = (nc.bytes_sent - self._prev_net[0]) / 1024
            recv_kb = (nc.bytes_recv - self._prev_net[1]) / 1024
            self._prev_net = (nc.bytes_sent, nc.bytes_recv)
        except Exception:
            sent_kb = recv_kb = 0.0
        n = self._nc_rows
        n[0].setText(f"↑ Send:   {sent_kb:.1f} KB/s"); n[0].setStyleSheet(f"color: {P['green']}; font-size: 12px;")
        n[1].setText(f"↓ Recv:   {recv_kb:.1f} KB/s"); n[1].setStyleSheet(f"color: {P['text2']}; font-size: 12px;")
        try: conns = len(psutil.net_connections(kind="inet"))
        except Exception: conns = 0
        n[2].setText(f"Conns:    {conns}"); n[2].setStyleSheet(f"color: {P['text3']}; font-size: 12px;")

        # graphs
        try:
            cpu = psutil.cpu_percent(interval=None); ram = psutil.virtual_memory().percent
        except Exception:
            cpu = ram = 0.0
        self._g_cpu.push(cpu); self._g_ram.push(ram)
        self._g_net_s.push(min(sent_kb, 1024)); self._g_net_r.push(min(recv_kb, 1024))

        # activity log auto-entries
        if running and self._tick % 2 == 0:
            _log_activity(f"Search engine: {phase}", "search")
        if self._tick % 5 == 0:
            _log_activity(f"System  CPU={cpu:.1f}%  RAM={ram:.1f}%", "sys")
        if self._tick % 10 == 0:
            _log_activity(f"Database  {total} jobs indexed", "db")

        # LLM new-call detection
        cur_count = len(log)
        if cur_count > self._prev_llm_count and log:
            newest = log[-1]
            _log_activity(f"LLM call: {newest['purpose'][:50]}  ({newest['duration_s']}s)", "info")
        self._prev_llm_count = cur_count
        self._lbl_llm_count.setText(f"{cur_count} calls")
        self._txt_llm.setPlainText("\n".join(
            f"[{e['ts']}] {e['purpose']:<40} {e['duration_s']}s" for e in log
        ))


        # AI decisions
        decisions = ai_score_engine.get_decisions()
        ok_count  = sum(1 for d in decisions if not d.get("error"))
        err_count = sum(1 for d in decisions if d.get("error"))
        count_txt = f"{ok_count} scored"
        if err_count:
            count_txt += f"  ·  {err_count} errors"
        self._lbl_ai_dec_count.setText(count_txt)
        dec_parts = []
        for d in reversed(decisions[-60:]):
            if d.get("error"):
                dec_parts.append(
                    f'<span style="color:{P["text3"]};">[{d["ts"]}]</span> '
                    f'<span style="color:{P["red"]};font-weight:700;">⚠ ERR</span> '
                    f'<span style="color:{P["text3"]};">{d["title"]}</span><br>'
                    f'<span style="color:{P["red"]};font-size:10px;">&nbsp;&nbsp;&nbsp;{d["reason"][:100]}</span>'
                )
                continue
            kept   = d["kept"]
            score  = d["score"]
            icon   = "✦" if kept else "✕"
            s_col  = P['green'] if score >= 60 else P['amber'] if score >= 30 else P['red']
            t_col  = P['text2'] if kept else P['text3']
            dec_parts.append(
                f'<span style="color:{P["text3"]};">[{d["ts"]}]</span> '
                f'<span style="color:{s_col};font-weight:700;">{icon} [{score:>3}]</span> '
                f'<span style="color:{t_col};">{d["title"]}</span>'
                f'<span style="color:{P["text3"]};"> · {d["company"]}</span><br>'
                f'<span style="color:{P["text3"]};font-size:10px;">&nbsp;&nbsp;&nbsp;{d["reason"][:90]}'
                f' &nbsp;<i>({d["time_s"]}s)</i></span>'
            )
        self._txt_ai_dec.setHtml("<br>".join(dec_parts) or
            f'<span style="color:{P["text3"]};">No AI decisions yet.</span>')
        sb2 = self._txt_ai_dec.verticalScrollBar(); sb2.setValue(sb2.minimum())

        # render activity log
        COLOR = {"search": P['indigo'], "sys": P['green'], "db": P['amber'], "info": P['text2']}
        parts = [
            f'<span style="color:{P["text3"]};">[{e["ts"]}]</span> '
            f'<span style="color:{COLOR.get(e["level"], P["text2"])};">{e["msg"]}</span>'
            for e in reversed(list(_ACTIVITY_LOG)[-60:])
        ]
        self._txt_activity.setHtml("<br>".join(parts))
        sb = self._txt_activity.verticalScrollBar(); sb.setValue(sb.minimum())
        self._lbl_events.setText(f"{len(_ACTIVITY_LOG)} events")


# ── helpers ───────────────────────────────────────────────────────────────────

def _profile_text(profile: dict) -> str:
    lines = [
        f"Name: {profile.get('name', '')}",
        f"Abschluss: {profile.get('degree', '')}",
        f"Hintergrund: {profile.get('background', '')}",
        f"Motivation: {profile.get('motivation', '')}",
    ]
    return "\n".join(l for l in lines if not l.endswith(": "))


def _set_palette(widget, bg: str, fg: str):
    pal = QPalette()
    pal.setColor(QPalette.ColorRole.Base, QColor(bg))
    pal.setColor(QPalette.ColorRole.Text, QColor(fg))
    widget.setPalette(pal)


def _make_app_icon() -> QIcon:
    px = QPixmap(256, 256); px.fill(QColor(0, 0, 0, 0))
    p = QPainter(px); p.setRenderHint(QPainter.RenderHint.Antialiasing)
    grad = QLinearGradient(0, 0, 0, 256)
    grad.setColorAt(0, QColor("#1e1e50")); grad.setColorAt(1, QColor("#08081a"))
    p.setBrush(QBrush(grad)); p.setPen(QPen(QColor("#6d6df5"), 10))
    p.drawRoundedRect(5, 5, 246, 246, 46, 46)
    p.setPen(QColor("#eeeef8")); p.setFont(QFont("Segoe UI Variable", 130, QFont.Weight.Bold))
    p.drawText(QRect(0, -10, 256, 256), Qt.AlignmentFlag.AlignCenter, "A")
    p.setPen(QPen(QColor("#6d6df5"), 14))
    p.drawLine(50, 198, 206, 198); p.drawLine(180, 174, 206, 198); p.drawLine(180, 222, 206, 198)
    p.end()
    return QIcon(px)


# ── entry point ───────────────────────────────────────────────────────────────

def run():
    import sys
    init_db()
    _log_activity("AutoApply started", "info")
    _log_activity("Database initialized", "db")
    _log_activity("UI loaded — waiting for commands", "sys")
    app = QApplication.instance() or QApplication(sys.argv)
    app.setStyleSheet(_qss())
    for name in ["Segoe UI Variable", "Segoe UI", "Inter"]:
        if name in QFontDatabase.families():
            app.setFont(QFont(name, 10)); break
    icon = _make_app_icon()
    app.setWindowIcon(icon)
    win = MainWindow(); win.setWindowIcon(icon); win.show()
    sys.exit(app.exec())
